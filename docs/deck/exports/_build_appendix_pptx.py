"""Build aivc_appendix_v1.pptx — 18-slide technical appendix.

Structure:
  Slide 01      Cover
  Slide 02      Section A divider · ARCHITECTURE DEPTH
  Slides 03-06  A1 / A2 / A3 / A4
  Slide 07      Section B divider · VALIDATION EVIDENCE
  Slides 08-10  B1 / B2 / B3
  Slide 11      Section C divider · QURIESEQ PHASE 1
  Slides 12-13  C1 / C2
  Slide 14      Section D divider · ROADMAP + BUDGET
  Slides 15-16  D1 / D2
  Slide 17      Section E divider · STRATEGIC HORIZON
  Slide 18      E1

Total: 1 cover + 5 dividers + 12 content = 18 slides.

Content slides are full-bleed PNG embeds (the PNG is already a self-contained
1920×1080 slide visual with title, body, source citation all rendered into the
SVG). The .pptx layer just sets the dark navy background and embeds the image.

Speaker notes for content slides are extracted from each content spec's
`## Speaker notes` section as plain-text Q&A (markdown stripped).

Run:  python3 docs/deck/exports/_build_appendix_pptx.py
"""
from __future__ import annotations
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

REPO = Path(__file__).resolve().parents[3]
CONTENT_DIR = REPO / "docs" / "deck" / "content"
DIAGRAMS_DIR = REPO / "docs" / "deck" / "assets" / "diagrams"
# v2: adds Section F (Competitive Positioning) + F1 content slide.
# v1 file preserved as historical artifact at aivc_appendix_v1.pptx.
OUTPUT = REPO / "docs" / "deck" / "exports" / "aivc_appendix_v2.pptx"

# 16:9 widescreen matching SVG viewBox 1920×1080 (13.333" × 7.5" at 144 dpi)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Palette (matches docs/deck/assets/color_palette.md)
BG_DARK     = RGBColor(0x07, 0x0A, 0x14)
CYAN        = RGBColor(0x26, 0xDD, 0xF9)
CYAN_HI     = RGBColor(0x00, 0xF2, 0xFF)
LAVENDER    = RGBColor(0xB4, 0x7D, 0xF0)
OK_GREEN    = RGBColor(0x4A, 0xDE, 0x80)
WARN_AMBER  = RGBColor(0xFB, 0xBF, 0x24)  # Section F accent — matches F1 SVG #FBBF24
TEXT_TITLE  = RGBColor(0xF7, 0xFA, 0xFF)
TEXT_BODY   = RGBColor(0xEA, 0xF6, 0xFF)
TEXT_MUTED  = RGBColor(0xA8, 0xB4, 0xC2)
TEXT_DIM    = RGBColor(0x94, 0xA3, 0xB8)

FONT_TITLE = "Inter"
FONT_BODY  = "Arial"

# Per-slide content map
# (slide_id, png_filename, content_spec_filename)
CONTENT_SLIDES = [
    ("A1", "A1_system_architecture_preview.png",         "A1_system_architecture.md"),
    ("A2", "A2_encoder_evidence_preview.png",            "A2_encoder_substrate.md"),
    ("A3", "A3_decomposed_readout_preview.png",          "A3_decomposed_readout.md"),
    ("A4", "A4_temporal_dynamics_preview.png",           "A4_temporal_neural_ode.md"),
    ("B1", "B1_three_datasets_methodology_preview.png",  "B1_methodology_rigor.md"),
    ("B2", "B2_adapter_verdict_preview.png",             "B2_encoder_probe_verdict.md"),
    ("B3", "B3_mechanism_pre_demo_preview.png",          "B3_synergy_pre_demo.md"),
    ("C1", "C1_phase1_experimental_design_preview.png",  "C1_phase1_design.md"),
    ("C2", "C2_btk_jak_demo_plan_preview.png",           "C2_btk_jak_demo.md"),
    ("D1", "D1_quarterly_roadmap_preview.png",           "D1_quarterly_roadmap.md"),
    ("D2", "D2_seed_allocation_preview.png",             "D2_seed_allocation.md"),
    ("E1", "E1_five_year_trajectory_preview.png",        "E1_five_year_trajectory.md"),
    ("F1", "F1_integrated_platform_preview.png",         "F1_competitive_positioning.md"),  # v2
]

SECTIONS = [
    {
        "letter": "A",
        "title": "Architecture Depth",
        "sub": "How the model works · what's frozen, what trains, what generalizes",
        "footer": "Slides A1 – A4",
        "accent": CYAN,
        "ids": ["A1", "A2", "A3", "A4"],
    },
    {
        "letter": "B",
        "title": "Validation Evidence",
        "sub": "How we test — three datasets, pre-registered evals, no cherry-picking",
        "footer": "Slides B1 – B3",
        "accent": OK_GREEN,
        "ids": ["B1", "B2", "B3"],
    },
    {
        "letter": "C",
        "title": "QurieSeq Phase 1",
        "sub": "The proprietary data engine · 5 donors × 5 timepoints × 4-arm × ~500K cells · BTK+JAK confirmed",
        "footer": "Slides C1 – C2",
        "accent": CYAN,
        "ids": ["C1", "C2"],
    },
    {
        "letter": "D",
        "title": "Roadmap + Budget",
        "sub": "11 quarters · 5 stages · 2 drug pipelines · $10M seed allocation",
        "footer": "Slides D1 – D2",
        "accent": LAVENDER,
        "ids": ["D1", "D2"],
    },
    {
        "letter": "E",
        "title": "Strategic Horizon",
        "sub": "From validated platform to first-in-class candidates · 2026 → 2031",
        "footer": "Slide E1",
        "accent": TEXT_BODY,
        "ids": ["E1"],
    },
    {
        # v2: new Section F appended after E. Amber accent matches F1 SVG.
        "letter": "F",
        "title": "Competitive Positioning",
        "sub": ("Why us · the closed-loop integrated platform · "
                "proprietary data, co-designed architecture, compounding over time"),
        "footer": "Slide F1",
        "accent": WARN_AMBER,
        "ids": ["F1"],
    },
]


# =============================================================================
# Speaker-note extraction
# =============================================================================
def _strip_markdown(text: str) -> str:
    """Convert markdown-flavored speaker note to plain text suitable for the
    PowerPoint notes slide. Preserves Q&A structure but strips emphasis."""
    # Strip bold: **text** → text
    text = re.sub(r"\*\*([^*]+)\*\*", r"\1", text)
    # Strip italic: *text* → text   (after bold so we don't break nested **bold**)
    text = re.sub(r"(?<!\*)\*([^*\n]+)\*(?!\*)", r"\1", text)
    # Strip blockquote markers: leading "> " becomes nothing
    text = re.sub(r"^>\s?", "", text, flags=re.MULTILINE)
    # Backticks: `code` → code
    text = re.sub(r"`([^`]+)`", r"\1", text)
    # Markdown links [text](url) → text
    text = re.sub(r"\[([^\]]+)\]\([^)]+\)", r"\1", text)
    # Collapse 3+ consecutive newlines to 2
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def extract_speaker_notes(spec_path: Path) -> str:
    """Pull the `## Speaker notes` section from a content spec and return
    plain-text Q&A. Returns empty string if no Speaker notes section."""
    if not spec_path.exists():
        return ""
    raw = spec_path.read_text(encoding="utf-8")
    # Find the speaker notes header (may have trailing parenthetical)
    m = re.search(r"^## Speaker notes[^\n]*\n", raw, re.MULTILINE)
    if not m:
        return ""
    start = m.end()
    # End at next "## " header or horizontal rule "---" on its own line
    end_match = re.search(r"^(## |\Z|---\s*$)", raw[start:], re.MULTILINE)
    end = start + end_match.start() if end_match else len(raw)
    section = raw[start:end].strip()
    return _strip_markdown(section)


# =============================================================================
# Slide-construction helpers
# =============================================================================
def set_slide_bg(slide, color: RGBColor) -> None:
    """Force a solid-fill background on a slide via OOXML manipulation.

    python-pptx exposes slide.background.fill, but for blank layouts the
    default fill chain often defers to the theme — setting via direct XML
    insertion is more reliable across PowerPoint versions.
    """
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_run_font(run, *, name: str, size_pt: int, bold: bool = False,
                  color: RGBColor = None, italic: bool = False) -> None:
    run.font.name = name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color


def _add_textbox(slide, *, left: float, top: float, width: float, height: float,
                 text: str, font_name: str, font_size: int, bold: bool = False,
                 italic: bool = False, color: RGBColor = TEXT_BODY,
                 align: PP_ALIGN = PP_ALIGN.LEFT,
                 letter_spacing_em: float = 0.0):
    """Add a text box at the specified inch coordinates."""
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    tf.word_wrap = True

    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    _set_run_font(run, name=font_name, size_pt=font_size, bold=bold,
                  italic=italic, color=color)

    if letter_spacing_em > 0:
        # python-pptx doesn't expose letter-spacing directly; use OOXML rPr/spc
        rPr = run._r.get_or_add_rPr()
        # spc value in 1/100ths of a point (hundredths)
        spc_val = int(letter_spacing_em * font_size * 100)
        rPr.set("spc", str(spc_val))
    return tb


def _add_divider_line(slide, *, left: float, top: float, width: float,
                      color: RGBColor = CYAN, thickness_pt: float = 1.5):
    """Thin horizontal rule. Implemented as a connector shape."""
    from pptx.enum.shapes import MSO_CONNECTOR
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(left), Inches(top),
        Inches(left + width), Inches(top),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(thickness_pt)
    return conn


# =============================================================================
# Slide builders
# =============================================================================
def add_cover_slide(prs):
    """Slide 01 — APPENDIX cover. Typography-only, no diagrams."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    # Eyebrow "APPENDIX" — cyan letter-spaced caps
    _add_textbox(
        slide, left=0.8, top=2.6, width=12, height=0.5,
        text="APPENDIX",
        font_name=FONT_TITLE, font_size=18, bold=True,
        color=CYAN, letter_spacing_em=0.3,
    )
    # Title
    _add_textbox(
        slide, left=0.8, top=3.1, width=12, height=1.0,
        text="AIVC GeneLink — Technical Appendix",
        font_name=FONT_TITLE, font_size=44, bold=True,
        color=TEXT_TITLE,
    )
    # Divider rule
    _add_divider_line(slide, left=0.8, top=4.25, width=4.5, color=CYAN)
    # Subtitle (section names)
    _add_textbox(
        slide, left=0.8, top=4.45, width=12, height=0.9,
        text=("Architecture Depth  ·  Validation Evidence  ·  QurieSeq Phase 1  ·  "
              "Roadmap + Budget  ·  Strategic Horizon"),
        font_name=FONT_BODY, font_size=18,
        color=TEXT_MUTED,
    )
    # Footer
    _add_textbox(
        slide, left=0.8, top=6.7, width=12, height=0.4,
        text="Quriegen · May 2026",
        font_name=FONT_TITLE, font_size=14, bold=True,
        color=CYAN, letter_spacing_em=0.2,
    )
    return slide


def add_section_divider(prs, section: dict):
    """Section divider slide. Same typography pattern as cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # Eyebrow "SECTION A"
    _add_textbox(
        slide, left=0.8, top=2.4, width=12, height=0.5,
        text=f"SECTION {section['letter']}",
        font_name=FONT_TITLE, font_size=16, bold=True,
        color=section["accent"], letter_spacing_em=0.3,
    )
    # Title
    _add_textbox(
        slide, left=0.8, top=2.9, width=12, height=1.4,
        text=section["title"],
        font_name=FONT_TITLE, font_size=56, bold=True,
        color=TEXT_TITLE,
    )
    # Divider rule
    _add_divider_line(slide, left=0.8, top=4.35, width=5.0, color=section["accent"])
    # Sub
    _add_textbox(
        slide, left=0.8, top=4.55, width=12, height=1.4,
        text=section["sub"],
        font_name=FONT_BODY, font_size=20,
        color=TEXT_MUTED,
    )
    # Footer (slides range)
    _add_textbox(
        slide, left=0.8, top=6.7, width=12, height=0.4,
        text=section["footer"],
        font_name=FONT_BODY, font_size=14,
        color=TEXT_DIM,
    )
    return slide


def add_content_slide(prs, png_path: Path, speaker_notes: str, slide_id: str):
    """Content slide — full-bleed PNG + speaker notes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    if not png_path.exists():
        raise FileNotFoundError(f"missing PNG for {slide_id}: {png_path}")

    # Full-bleed: PNG fills the entire slide (1920×1080 ratio matches 13.333" × 7.5")
    slide.shapes.add_picture(
        str(png_path),
        left=0, top=0,
        width=SLIDE_WIDTH, height=SLIDE_HEIGHT,
    )

    # Speaker notes
    if speaker_notes:
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        # Clear default and set
        tf.text = speaker_notes
    return slide


# =============================================================================
# Main
# =============================================================================
def main() -> int:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # --- Slide 01: Cover ---
    add_cover_slide(prs)

    # --- Sections + content slides ---
    for section in SECTIONS:
        # Section divider
        add_section_divider(prs, section)
        # Content slides in this section
        for slide_id in section["ids"]:
            cs = next(c for c in CONTENT_SLIDES if c[0] == slide_id)
            _, png_name, spec_name = cs
            png_path = DIAGRAMS_DIR / png_name
            spec_path = CONTENT_DIR / spec_name
            speaker_notes = extract_speaker_notes(spec_path)
            add_content_slide(prs, png_path, speaker_notes, slide_id)

    # --- Save ---
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)

    # Audit
    n_slides = len(prs.slides)
    file_size = OUTPUT.stat().st_size
    print(f"Built {OUTPUT}")
    print(f"  slides:   {n_slides}")
    print(f"  size:     {file_size / 1024:.1f} KB ({file_size} bytes)")
    expected_total = 1 + len(SECTIONS) + sum(len(s["ids"]) for s in SECTIONS)
    if n_slides != expected_total:
        print(f"  WARN: expected {expected_total} slides, got {n_slides}", file=sys.stderr)
        return 1
    if file_size > 10 * 1024 * 1024:
        print(f"  WARN: file exceeds 10MB budget", file=sys.stderr)
    return 0


if __name__ == "__main__":
    sys.exit(main())
