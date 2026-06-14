"""AIVC Seed Round Plan — wet + dry lab expansion (Path A native shapes).

Replicates the source slide structure ("Seed round plan: expand wet and dry
lab capabilities" / SOLUTION · AI VIRTUAL CELL MODEL) with the v3 design
system applied. Three biological scales (Submolecular / Molecular / Tissue)
inside the platform section, two enabling pillars (Data Moat + Wet Lab)
below.

Color discipline (scale-based, internally consistent with v3 palette):
  - Submolecular = cyan       (smallest scale, foundational)
  - Molecular    = lavender   (mid scale)
  - Tissue       = amber      (largest scale)
  - Enablers     = green      (data + lab infrastructure)

Output: docs/deck/exports/aivc_seed_plan_v1.pptx (single slide)

Source-slide cleanups applied:
  1. QuRIE-Flow/QuRIE-Perturb consolidated under Tissue column only
     (source had a dashed-bordered orphan duplicate under Submolecular).
  2. "QURIE-Recon" normalized to "QuRIE-Recon" for naming consistency with
     QuRIE-seq / QuRIE-PerturbSeq family (flag for Ash override).
  3. Removed the blue rounded container (visual clutter); platform scope
     conveyed by typography + accent rules + a thin enabler divider.

Run:
  python3 docs/deck/investor_4slide/_build_seed_plan_slide.py
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Locked palette
# ---------------------------------------------------------------------------
BG_DARK         = RGBColor(0x0A, 0x0E, 0x1A)
FG_PRIMARY      = RGBColor(0xFF, 0xFF, 0xFF)
FG_SECONDARY    = RGBColor(0xA0, 0xAF, 0xC8)
FG_MUTED        = RGBColor(0x60, 0x70, 0x88)
ACCENT_CYAN     = RGBColor(0x26, 0xDD, 0xF9)
ACCENT_LAVENDER = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_GREEN    = RGBColor(0x4A, 0xDE, 0x80)
BORDER_SUBTLE   = RGBColor(0x2D, 0x3A, 0x57)

FONT = "Calibri"

REPO = Path(__file__).resolve().parents[3]
OUT_PATH = REPO / "docs" / "deck" / "exports" / "aivc_seed_plan_v1.pptx"


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------
def _zero_margins(tf):
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)


def _spc(run, hundredths_pt):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(hundredths_pt)))


def set_slide_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_text(slide, left, top, width, height, text,
             size=14, bold=False, italic=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, letter_spacing=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    _zero_margins(tf)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    if color is not None:
        run.font.color.rgb = color
    if letter_spacing is not None:
        _spc(run, letter_spacing)
    return box


def add_multiline(slide, left, top, width, height, lines,
                  size=13, color=None, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.35,
                  letter_spacing=None, bold=False, italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    _zero_margins(tf)
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        spec = item if isinstance(item, dict) else {"text": item}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", align)
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = spec["text"]
        run.font.size = Pt(spec.get("size", size))
        run.font.bold = spec.get("bold", bold)
        run.font.italic = spec.get("italic", italic)
        run.font.name = FONT
        col = spec.get("color", color)
        if col is not None:
            run.font.color.rgb = col
        ls = spec.get("letter_spacing", letter_spacing)
        if ls is not None:
            _spc(run, ls)
    return box


def add_accent_rule(slide, x, y, length, color, thickness_in=0.05):
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y),
        Inches(length), Inches(thickness_in),
    )
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_divider(slide, x1, y, x2, color, width=0.75):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y), Inches(x2), Inches(y),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# Layout constants
# ---------------------------------------------------------------------------
# Top zone
CRUMB_Y         = 0.35
TITLE_Y         = 0.62
SUBTITLE_Y      = 1.40
TITLE_DIVIDER_Y = 1.88

# Scale section (top half)
SCALE_HEADER_Y      = 2.02
ACCENT_RULE_Y       = 2.34
MODULES_START_Y     = 2.55
MODULE_HEIGHT       = 1.05
INTRA_MODULE_GAP    = 0.05

# Per-module internal layout offsets (relative to module top)
MOD_CAT_OFFSET   = 0.00
MOD_CAT_HEIGHT   = 0.20
MOD_NAME_OFFSET  = 0.22
MOD_NAME_HEIGHT  = 0.32
MOD_DESC_OFFSET  = 0.55
MOD_DESC_HEIGHT  = 0.52

# Scale column layout (3 columns)
COLUMN_LEFTS    = [0.60, 4.85, 9.05]
COLUMN_WIDTH    = 3.65
COLUMN_CENTERS  = [L + COLUMN_WIDTH / 2 for L in COLUMN_LEFTS]
ACCENT_RULE_W   = 1.50

# Enabler section divider + label
ENABLER_DIVIDER_Y = 5.95
ENABLER_LABEL_Y   = 6.05

# Enabler row (2 boxes split horizontally)
ENABLER_LEFTS   = [0.60, 6.95]
ENABLER_WIDTH   = 5.78
ENABLER_HEADER_Y   = 6.05   # category label
ENABLER_NAME_Y     = 6.30   # product name (bold green)
ENABLER_DESC_Y     = 6.75   # description (white body)

# Footer
FOOTER_Y = 7.18


# ---------------------------------------------------------------------------
# Module renderer
# ---------------------------------------------------------------------------
def render_module(slide, left, top, width, color, category, name,
                  description, name_size=18):
    """Render one product module within a scale column.

    Structure:
      - Intelligence category label (italic, FG_SECONDARY, 11pt)
      - Product name (bold, scale color, 18pt by default)
      - Description (regular, FG_PRIMARY, 12pt, 2-3 line wrap)
    """
    # Category label
    add_text(slide, left, top + MOD_CAT_OFFSET, width, MOD_CAT_HEIGHT,
             category, size=11, italic=True, color=FG_SECONDARY,
             align=PP_ALIGN.LEFT, line_spacing=1.1)

    # Product name
    add_text(slide, left, top + MOD_NAME_OFFSET, width, MOD_NAME_HEIGHT,
             name, size=name_size, bold=True, color=color,
             align=PP_ALIGN.LEFT, line_spacing=1.0)

    # Description (tight line spacing to fit 2-line descriptions reliably)
    add_text(slide, left, top + MOD_DESC_OFFSET, width, MOD_DESC_HEIGHT,
             description, size=12, color=FG_PRIMARY,
             align=PP_ALIGN.LEFT, line_spacing=1.20)


def render_scale_column(slide, idx, color, scale_label, modules):
    """Render a complete scale column (header + accent rule + stacked modules).

    `modules` is a list of (category, name, description) tuples.
    """
    left = COLUMN_LEFTS[idx]

    # Scale header (uppercase, letter-spaced, scale color)
    add_text(slide, left, SCALE_HEADER_Y, COLUMN_WIDTH, 0.30,
             scale_label, size=14, bold=True, color=color,
             align=PP_ALIGN.LEFT, line_spacing=1.0, letter_spacing=200)

    # Accent rule (1.5" wide, left-aligned)
    add_accent_rule(slide, left, ACCENT_RULE_Y, ACCENT_RULE_W, color,
                    thickness_in=0.05)

    # Modules stacked vertically
    for i, (category, name, description) in enumerate(modules):
        mod_top = MODULES_START_Y + i * (MODULE_HEIGHT + INTRA_MODULE_GAP)
        render_module(slide, left, mod_top, COLUMN_WIDTH, color,
                      category, name, description)


def render_enabler(slide, idx, label, name, description):
    """Render an enabler box (Data Moat / Wet Lab) — text-only, no border."""
    left = ENABLER_LEFTS[idx]
    width = ENABLER_WIDTH

    # Category label
    add_text(slide, left, ENABLER_HEADER_Y, width, 0.22,
             label, size=12, italic=True, color=FG_SECONDARY,
             align=PP_ALIGN.LEFT, line_spacing=1.0, letter_spacing=100)

    # Product name (bold green)
    add_text(slide, left, ENABLER_NAME_Y, width, 0.40,
             name, size=22, bold=True, color=ACCENT_GREEN,
             align=PP_ALIGN.LEFT, line_spacing=1.0)

    # Description
    add_text(slide, left, ENABLER_DESC_Y, width, 0.70,
             description, size=12, color=FG_PRIMARY,
             align=PP_ALIGN.LEFT, line_spacing=1.30)


# ---------------------------------------------------------------------------
# Slide build
# ---------------------------------------------------------------------------
def build_seed_plan_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    # ── Top crumb header ──────────────────────────────────────────────────
    add_text(slide, 0.6, CRUMB_Y, 12.13, 0.25,
             "SOLUTION  ·  AI VIRTUAL CELL MODEL",
             size=11, bold=True, color=ACCENT_CYAN,
             align=PP_ALIGN.LEFT, line_spacing=1.0, letter_spacing=200)

    # ── Title + subtitle ──────────────────────────────────────────────────
    add_text(slide, 0.6, TITLE_Y, 12.13, 0.80,
             "Seed Round Plan",
             size=44, bold=True, color=FG_PRIMARY,
             align=PP_ALIGN.LEFT, line_spacing=1.0)
    add_text(slide, 0.6, SUBTITLE_Y, 12.13, 0.40,
             "Expanding wet and dry lab capabilities across biological scales",
             size=18, italic=True, color=FG_SECONDARY,
             align=PP_ALIGN.LEFT, line_spacing=1.1)
    add_divider(slide, 0.6, TITLE_DIVIDER_Y, 12.733, BORDER_SUBTLE, width=0.75)

    # ── Submolecular column (cyan, 1 module) ──────────────────────────────
    render_scale_column(slide, 0, ACCENT_CYAN,
        scale_label="SUBMOLECULAR",
        modules=[
            (
                "Mechanistic Intelligence",
                "Qurie-Kinetica",
                "Mechanistic modeling of signaling, reaction kinetics, "
                "and cell dynamics.",
            ),
        ],
    )

    # ── Molecular column (lavender, 3 modules) ────────────────────────────
    render_scale_column(slide, 1, ACCENT_LAVENDER,
        scale_label="MOLECULAR",
        modules=[
            (
                "Physical Cell Intelligence",
                "Qurie-PhysioMap",
                "Imaging, morphology (size, shape), protein-localization, "
                "and mechanics readouts.",
            ),
            (
                "Cell State Intelligence",
                "QuRIE-seq",
                "Multimodal single-cell profiling of RNA, proteins, "
                "and phosphoproteins.",
            ),
            (
                "Causal Cell Intelligence",
                "QuRIE-PerturbSeq",
                "CRISPR perturbation linked to RNA, proteins, "
                "phosphoproteins, and VDJ.",
            ),
        ],
    )

    # ── Tissue column (amber, 2 modules) ──────────────────────────────────
    render_scale_column(slide, 2, ACCENT_AMBER,
        scale_label="TISSUE",
        modules=[
            (
                "Tissue Intelligence",
                "QuRIE-Recon",
                "Tumor–immune tissue reconstruction. "
                "Therapy + tissue-aware biomarker prediction.",
            ),
            (
                "Tissue Intelligence",
                "QuRIE-Flow / QuRIE-Perturb",
                "AI software for 2D/3D tissue reconstruction from "
                "dissociated single cells.",
            ),
        ],
    )

    # ── Enabler section divider + label ───────────────────────────────────
    add_divider(slide, 0.6, ENABLER_DIVIDER_Y, 12.733, BORDER_SUBTLE, width=0.75)

    # ── Two enabler boxes (Data Moat + Wet Lab, both green) ───────────────
    render_enabler(
        slide, 0,
        label="DATA MOAT  ·  Bioinformatic Intelligence",
        name="Qurie-PBMC Atlas",
        description=(
            "Largest deeply profiled immune atlas — enables the first dynamic "
            "immune clock and benchmarks immune responses across perturbation, "
            "aging, and disease."
        ),
    )
    render_enabler(
        slide, 1,
        label="WET LAB  ·  Lab-in-the-loop",
        name="Qurie-Lab",
        description=(
            "Laboratory capabilities for data generation, perturbation "
            "biology, and model validation."
        ),
    )

    add_speaker_notes(slide, _SPEAKER_NOTES)
    return slide


# ---------------------------------------------------------------------------
# Speaker notes
# ---------------------------------------------------------------------------
_SPEAKER_NOTES = """\
Seed Round Plan — Wet + Dry Lab Expansion (90-120s walk-through)

OPENING (10s):
"This slide shows what the seed round funds: an integrated platform spanning
three biological scales, supported by a proprietary data moat and our own
wet lab. Every module is a distinct deliverable; together they compose into
the AIVC virtual cell model."

SUBMOLECULAR SCALE (15s):
"At the submolecular scale: Qurie-Kinetica — the mechanistic intelligence
layer. Models signaling dynamics, reaction kinetics, and cell dynamics
under physics-informed priors. This is where we extrapolate beyond observed
perturbations into resistance evolution and combination effects."

MOLECULAR SCALE (30s):
"At the molecular scale: three complementary modules.

Qurie-PhysioMap — Physical Cell Intelligence. Imaging-derived morphology,
protein localization, and mechanics readouts.

QuRIE-seq — Cell State Intelligence. Our proprietary multimodal single-cell
assay measuring RNA, proteins, and phosphoproteins from the same cell.

QuRIE-PerturbSeq — Causal Cell Intelligence. CRISPR perturbation linked
to RNA, proteins, phosphoproteins, and VDJ — the causal layer in molecular
intelligence."

TISSUE SCALE (20s):
"At the tissue scale: two modules.

QuRIE-Recon — tissue reconstruction and tumor–immune interaction modeling.
Drives therapy prediction and tissue-aware biomarker discovery.

QuRIE-Flow / QuRIE-Perturb — the AI software stack for 2D/3D tissue
reconstruction from dissociated single cells. This is how we get from
single-cell measurements back to tissue context."

ENABLING PILLARS (20s):
"Two infrastructure pillars support all six modules.

Qurie-PBMC Atlas — our data moat. Largest deeply profiled immune atlas.
Enables the first dynamic immune clock and benchmarks immune responses
across perturbation, aging, and disease.

Qurie-Lab — our wet lab. Lab-in-the-loop capabilities for data generation,
perturbation biology, and model validation. The proprietary input layer
that makes the rest of the platform defensible."

CLOSE (10s):
"Seven products plus two pillars — but one platform. Every module operates
on the same shared biological state, generated by Qurie-Lab and
benchmarked against Qurie-PBMC Atlas. That's how this seed round
compounds: each capability strengthens every other."

EXPECTED DILIGENCE QUESTIONS:

Q: "Why three biological scales — isn't this scope creep?"
A: Scale is where biological context lives. Submolecular for mechanism,
molecular for measurement, tissue for clinical relevance. Skipping any
scale breaks the chain from data to therapy.

Q: "Does the seed fund all seven modules in parallel?"
A: No. Seed funds Qurie-Lab + Qurie-PBMC Atlas (infrastructure pillars)
plus Qurie-Kinetica + QuRIE-seq + QuRIE-PerturbSeq (the core triad).
PhysioMap + Recon + Flow/Perturb are Series A roadmap with seed-stage
proof-of-concept work.

Q: "How is this different from Insilico / Recursion / CytoReason?"
A: Three things they don't have together: proprietary multimodal single-cell
perturbation data (QuRIE-seq), mechanistic priors as differentiable
constraints (Qurie-Kinetica), and tissue-scale reconstruction from
dissociated cells. The integration is the moat, not any single module.

Q: "Naming — why the QuRIE / Qurie inconsistency?"
A: Brand convention in flight. Models = 'Qurie-X' (Qurie-Kinetica,
Qurie-PhysioMap, Qurie-Lab). Assays = 'QuRIE-X' (QuRIE-seq,
QuRIE-PerturbSeq). Will normalize before Series A pitch.
"""


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_seed_plan_slide(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    size_kb = OUT_PATH.stat().st_size / 1024
    print(f"Saved: {OUT_PATH}  ({size_kb:.1f} KB)")


if __name__ == "__main__":
    main()
