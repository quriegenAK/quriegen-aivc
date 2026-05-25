"""Investor 4-Slide Deck builder — Path A native python-pptx shapes only.

Spec: docs/deck/prompts/investor_4slide_build_spec.md

Critical constraint:
  NO flattened images for diagrams or text. All visual elements must be
  native PowerPoint shapes (rectangles, ovals, connectors, text boxes)
  editable by reviewers in PowerPoint Mac.

Output: docs/deck/exports/aivc_investor_4slide_v1.pptx

Run:
  python3 docs/deck/investor_4slide/_build_investor_deck.py
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Color palette (single source of truth, per spec)
# ---------------------------------------------------------------------------
BG_DARK         = RGBColor(0x0A, 0x0E, 0x1A)
FG_PRIMARY      = RGBColor(0xFF, 0xFF, 0xFF)
FG_SECONDARY    = RGBColor(0xA0, 0xAF, 0xC8)
ACCENT_CYAN     = RGBColor(0x26, 0xDD, 0xF9)
ACCENT_LAVENDER = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_GREEN    = RGBColor(0x4A, 0xDE, 0x80)
BORDER_SUBTLE   = RGBColor(0x2D, 0x3A, 0x57)

FONT = "Calibri"

REPO = Path(__file__).resolve().parents[3]
OUT_PATH = REPO / "docs" / "deck" / "exports" / "aivc_investor_4slide_v1.pptx"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def set_slide_bg(slide, rgb):
    """Solid background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=14,
    bold=False,
    italic=False,
    color=None,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.15,
):
    """Add a text box with a single styled run."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    if color is not None:
        run.font.color.rgb = color
    return box


def add_multiline_text(
    slide,
    left,
    top,
    width,
    height,
    lines,
    size=14,
    bold=False,
    italic=False,
    color=None,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.25,
):
    """Add a multi-line text box. `lines` is a list of strings OR dicts
    {text, size?, bold?, italic?, color?}. Each becomes its own paragraph.
    """
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        if isinstance(item, str):
            spec = {"text": item}
        else:
            spec = dict(item)
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", align)
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = spec["text"]
        run.font.size = Pt(spec.get("size", size))
        run.font.bold = spec.get("bold", bold)
        run.font.italic = spec.get("italic", italic)
        run.font.name = FONT
        col = spec.get("color", color)
        if col is not None:
            run.font.color.rgb = col
    return box


def add_card(
    slide,
    left,
    top,
    width,
    height,
    fill_color=None,
    border_color=BORDER_SUBTLE,
    border_width=1.0,
    corner=0.06,
):
    """Rounded rectangle. Pass fill_color=None for transparent fill."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.adjustments[0] = corner
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color is not None:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=1.0):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    shape.shadow.inherit = False
    return shape


def _set_fill_transparency(shape, alpha_percent):
    """Apply solid fill alpha (0-100) to a shape's solidFill XML element."""
    sppr = shape.fill._xPr.find(qn("p:spPr")) if False else shape.fill._xPr
    # python-pptx exposes fill xml via shape.fill._xPr; the solidFill element
    # already exists when fill.solid() was called. We modify its srgbClr child.
    spPr = shape.fill._xPr
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is None:
        return
    srgb = solidFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    # Remove existing alpha to make idempotent
    for existing in srgb.findall(qn("a:alpha")):
        srgb.remove(existing)
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", str(int(alpha_percent * 1000)))


def add_card_translucent(slide, left, top, width, height, fill_color, alpha_pct, border_color, border_width=1.5, corner=0.06):
    """Rounded rectangle with translucent fill (alpha 0-100)."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.adjustments[0] = corner
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    _set_fill_transparency(shape, alpha_pct)
    if border_color is not None:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_arrow(slide, x1, y1, x2, y2, color, width=1.5, head=True):
    """Straight connector (line). Optionally add a triangle arrowhead at the tail end."""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    if head:
        line_elem = connector.line._get_or_add_ln()
        # Remove existing tailEnd to make idempotent
        for existing in line_elem.findall(qn("a:tailEnd")):
            line_elem.remove(existing)
        tail = etree.SubElement(line_elem, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", "med")
        tail.set("len", "med")
    return connector


def add_oval(slide, cx, cy, diameter, fill_color, border_color, border_width=1.0):
    r = diameter / 2.0
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - r),
        Inches(cy - r),
        Inches(diameter),
        Inches(diameter),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(border_width)
    shape.shadow.inherit = False
    return shape


def add_hline(slide, left, top, width, color, weight=1.0):
    """Horizontal divider line."""
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(left),
        Inches(top),
        Inches(left + width),
        Inches(top),
    )
    connector.line.color.rgb = color
    connector.line.width = Pt(weight)
    return connector


def add_title_block(slide, title, subtitle):
    """Standard slide title + subtitle at the top."""
    add_text(slide, 0.5, 0.35, 12.5, 0.7, title, size=32, bold=True, color=FG_PRIMARY)
    add_text(slide, 0.5, 0.95, 12.5, 0.45, subtitle, size=16, italic=True, color=FG_SECONDARY)
    add_hline(slide, 0.5, 1.5, 12.333, BORDER_SUBTLE, weight=0.75)


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, BG_DARK)
    return slide


# ---------------------------------------------------------------------------
# Slide 1 — AIVC Platform Evolution
# ---------------------------------------------------------------------------
def _modality_chip(slide, left, top, label, color):
    """Small geometric icon (filled circle) + label below for modality chips."""
    # Filled circle as icon
    diameter = 0.32
    add_oval(slide, left + diameter / 2, top + diameter / 2, diameter, color, color, border_width=0.5)
    add_text(
        slide,
        left - 0.4,
        top + 0.35,
        diameter + 0.8,
        0.25,
        label,
        size=10,
        color=FG_SECONDARY,
        align=PP_ALIGN.CENTER,
    )


def _phase_panel(slide, left, header_label, header_color, header_subtitle, when_text,
                 body_lines, italic_caption, visual_fn, panel_top=1.7, panel_width=4.05, panel_height=5.0):
    """Render one phase panel (Now/Phase 1/Phase 2). visual_fn(left, top) renders the diagram."""
    # Panel border card
    add_card(slide, left, panel_top, panel_width, panel_height,
             fill_color=None, border_color=BORDER_SUBTLE, border_width=1.0, corner=0.04)

    # Header label (small caps)
    add_text(slide, left + 0.25, panel_top + 0.18, panel_width - 0.5, 0.3,
             header_label, size=12, bold=True, color=header_color)

    # Header title
    add_text(slide, left + 0.25, panel_top + 0.5, panel_width - 0.5, 0.5,
             header_subtitle, size=18, bold=True, color=FG_PRIMARY)

    # Optional when text
    next_top = panel_top + 0.95
    if when_text:
        add_text(slide, left + 0.25, next_top, panel_width - 0.5, 0.3,
                 when_text, size=11, color=header_color)
        next_top += 0.35

    # Body — multi-line stack
    add_multiline_text(slide, left + 0.25, next_top, panel_width - 0.5, 2.0,
                       body_lines, size=13, color=FG_PRIMARY, line_spacing=1.35)

    # Visual element block
    visual_fn(left, panel_top + 3.45)

    # Italic caption at bottom of panel
    add_text(slide, left + 0.25, panel_top + panel_height - 0.7,
             panel_width - 0.5, 0.6, italic_caption,
             size=11, italic=True, color=FG_SECONDARY, align=PP_ALIGN.LEFT)


def _now_visual(slide, panel_left, top):
    """3 modality chips (RNA, ATAC, Protein) horizontally."""
    chips_y = top + 0.25
    chip_x_positions = [panel_left + 0.65, panel_left + 1.85, panel_left + 3.05]
    labels = ["RNA", "ATAC", "Protein"]
    for x, label in zip(chip_x_positions, labels):
        _modality_chip(slide, x, chips_y, label, ACCENT_CYAN)


def _phase1_visual(slide, panel_left, top):
    """3 modality boxes feeding into central encoder hexagon."""
    # 3 small rounded modality boxes top row
    box_w, box_h = 0.85, 0.4
    box_y = top + 0.05
    spacing = 1.05
    base_x = panel_left + (4.05 - (3 * box_w + 2 * (spacing - box_w))) / 2
    # Actually simpler: distribute 3 evenly
    panel_w = 4.05
    margin = 0.45
    usable = panel_w - 2 * margin
    gap = (usable - 3 * box_w) / 2
    xs = [panel_left + margin + i * (box_w + gap) for i in range(3)]
    labels = ["RNA", "Prot", "Phos"]
    for x, label in zip(xs, labels):
        s = add_card(slide, x, box_y, box_w, box_h,
                     fill_color=None, border_color=ACCENT_LAVENDER,
                     border_width=1.25, corner=0.18)
        add_text(slide, x, box_y + 0.04, box_w, box_h - 0.05,
                 label, size=10, bold=True, color=FG_PRIMARY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Central encoder hexagon
    hex_w, hex_h = 1.65, 0.65
    hex_x = panel_left + (panel_w - hex_w) / 2
    hex_y = box_y + box_h + 0.42
    hex_shape = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Inches(hex_x), Inches(hex_y), Inches(hex_w), Inches(hex_h),
    )
    hex_shape.fill.solid()
    hex_shape.fill.fore_color.rgb = ACCENT_LAVENDER
    _set_fill_transparency(hex_shape, 22)
    hex_shape.line.color.rgb = ACCENT_LAVENDER
    hex_shape.line.width = Pt(1.5)
    hex_shape.shadow.inherit = False
    add_text(slide, hex_x, hex_y + 0.08, hex_w, hex_h - 0.1,
             "Encoder", size=10, bold=True, color=FG_PRIMARY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 3 short arrows from each box to encoder
    for x in xs:
        x_center = x + box_w / 2
        add_arrow(slide, x_center, box_y + box_h, x_center, hex_y + 0.02,
                  ACCENT_LAVENDER, width=1.0, head=True)


def _phase2_visual(slide, panel_left, top):
    """5 modality boxes feeding into encoder + CRISPR/VDJ side arrows."""
    panel_w = 4.05
    box_w, box_h = 0.6, 0.36
    box_y = top + 0.05
    margin = 0.18
    usable = panel_w - 2 * margin
    gap = (usable - 5 * box_w) / 4
    xs = [panel_left + margin + i * (box_w + gap) for i in range(5)]
    labels = ["RNA", "Epi", "Prot", "Phos", "VDJ"]
    colors = [ACCENT_CYAN, ACCENT_CYAN, ACCENT_CYAN, ACCENT_LAVENDER, ACCENT_AMBER]
    for x, label, color in zip(xs, labels, colors):
        s = add_card(slide, x, box_y, box_w, box_h,
                     fill_color=None, border_color=color,
                     border_width=1.25, corner=0.18)
        add_text(slide, x, box_y + 0.025, box_w, box_h - 0.05,
                 label, size=9, bold=True, color=FG_PRIMARY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Encoder hexagon, slightly larger than Phase 1
    hex_w, hex_h = 1.85, 0.7
    hex_x = panel_left + (panel_w - hex_w) / 2
    hex_y = box_y + box_h + 0.42
    hex_shape = slide.shapes.add_shape(
        MSO_SHAPE.HEXAGON,
        Inches(hex_x), Inches(hex_y), Inches(hex_w), Inches(hex_h),
    )
    hex_shape.fill.solid()
    hex_shape.fill.fore_color.rgb = ACCENT_AMBER
    _set_fill_transparency(hex_shape, 22)
    hex_shape.line.color.rgb = ACCENT_AMBER
    hex_shape.line.width = Pt(1.5)
    hex_shape.shadow.inherit = False
    add_text(slide, hex_x, hex_y + 0.08, hex_w, hex_h - 0.1,
             "Encoder", size=10, bold=True, color=FG_PRIMARY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 5 short arrows from each box to encoder
    for x, color in zip(xs, colors):
        x_center = x + box_w / 2
        add_arrow(slide, x_center, box_y + box_h, x_center, hex_y + 0.02,
                  color, width=1.0, head=True)


def build_slide1_evolution(prs):
    slide = blank_slide(prs)
    add_title_block(slide,
                    "AIVC Platform Evolution",
                    "From public benchmarking to scalable causal biological intelligence")

    # 3 panel layout
    panel_w = 4.05
    panel_top = 1.75
    panel_height = 5.0
    gap = 0.15
    total = 3 * panel_w + 2 * gap
    start = (13.333 - total) / 2  # center the 3 panels
    panel_lefts = [start + i * (panel_w + gap) for i in range(3)]

    # Panel 1 — NOW
    _phase_panel(
        slide, panel_lefts[0],
        header_label="NOW",
        header_color=ACCENT_CYAN,
        header_subtitle="Foundation & Benchmarking",
        when_text=None,
        body_lines=[
            "Public multimodal datasets",
            "3 reference papers",
            {"text": "DOGMA-seq · Calderon · Mimitou CRISPR", "size": 11, "color": FG_SECONDARY, "italic": True},
            "Multi-omics encoder pretraining",
            {"text": "73% cross-corpus accuracy", "size": 14, "bold": True, "color": ACCENT_CYAN},
        ],
        italic_caption='"Foundational biological representation system, built on validated public data."',
        visual_fn=lambda L, T: _now_visual(slide, L, T),
        panel_top=panel_top, panel_width=panel_w, panel_height=panel_height,
    )

    # Panel 2 — PHASE 1
    _phase_panel(
        slide, panel_lefts[1],
        header_label="PHASE 1",
        header_color=ACCENT_LAVENDER,
        header_subtitle="Controlled Perturbation Learning",
        when_text="Q3 2026",
        body_lines=[
            "QuRIE-seq · proprietary multi-omics",
            "3 modalities directly measured:",
            {"text": "RNA · Proteins · Phosphoproteins", "size": 12, "color": ACCENT_LAVENDER, "bold": True},
            "PBMCs · 5 donors · 5 timepoints",
            "5 stimuli · 10 inhibitors",
            {"text": "BTK + JAK combo", "size": 13, "bold": True, "color": ACCENT_LAVENDER},
        ],
        italic_caption='"Learning causal biological responses under controlled perturbations."',
        visual_fn=lambda L, T: _phase1_visual(slide, L, T),
        panel_top=panel_top, panel_width=panel_w, panel_height=panel_height,
    )

    # Panel 3 — PHASE 2
    _phase_panel(
        slide, panel_lefts[2],
        header_label="PHASE 2",
        header_color=ACCENT_AMBER,
        header_subtitle="Scalable Causal Discovery",
        when_text="2027",
        body_lines=[
            "QuRIE-seq + CRISPR + VDJ",
            "5 modalities:",
            {"text": "RNA · Epi · Prot · Phospho · VDJ", "size": 11, "color": ACCENT_AMBER, "bold": True},
            "20–25 donors",
            "Soft perturbations: 30 stimuli + inhibitors",
            "Hard perturbations: CRISPR library",
        ],
        italic_caption='"Scaling toward large multimodal causal biological intelligence."',
        visual_fn=lambda L, T: _phase2_visual(slide, L, T),
        panel_top=panel_top, panel_width=panel_w, panel_height=panel_height,
    )

    # Inter-panel arrows (Now → Phase 1 → Phase 2)
    arrow_y = panel_top + panel_height / 2
    for i in range(2):
        x_from = panel_lefts[i] + panel_w + 0.005
        x_to = panel_lefts[i + 1] - 0.005
        # Tiny arrow at gap midpoint
        add_arrow(slide, x_from, arrow_y, x_to, arrow_y,
                  FG_SECONDARY, width=1.25, head=True)

    # Footer
    footer_top = 6.95
    add_hline(slide, 0.5, footer_top, 12.333, BORDER_SUBTLE, weight=0.75)
    add_text(slide, 0.5, footer_top + 0.12, 12.333, 0.4,
             "PHASE 3 — Continuation of data generation at scale + therapeutic pipeline",
             size=12, italic=True, color=FG_SECONDARY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 2 — Causal Model + Validation
# ---------------------------------------------------------------------------
def _draw_network(slide, cx, cy, radius, n_nodes, node_diameter, node_fill, node_border,
                  edges, edge_color, edge_width, directed=False, perturbation_node=None,
                  arrow_specs=None, angle_offset=0.0):
    """Draw a small network. nodes laid out roughly circularly around (cx, cy).

    edges: list of (i, j) index pairs.
    arrow_specs: optional dict {(i,j): (color, width, head_bool)} overrides per-edge styling.
    perturbation_node: if not None, the index gets accent fill + lightning indicator.
    """
    import math
    # Compute positions (slight inner randomness for organic look — but deterministic)
    positions = []
    # Mix: place ~half on outer ring, ~half on inner ring
    for i in range(n_nodes):
        angle = (2 * math.pi * i / n_nodes) - math.pi / 2 + angle_offset
        # Alternate inner / outer ring
        r = radius if i % 2 == 0 else radius * 0.62
        x = cx + r * math.cos(angle)
        y = cy + r * math.sin(angle)
        positions.append((x, y))

    # Draw edges first (so nodes overlap on top)
    for (i, j) in edges:
        x1, y1 = positions[i]
        x2, y2 = positions[j]
        if arrow_specs and (i, j) in arrow_specs:
            c, w, head = arrow_specs[(i, j)]
        else:
            c, w, head = edge_color, edge_width, directed
        # Offset endpoints by node radius so arrow tips hit edge, not center
        # Compute unit vector
        dx = x2 - x1
        dy = y2 - y1
        dist = (dx ** 2 + dy ** 2) ** 0.5
        if dist > 0:
            ux, uy = dx / dist, dy / dist
            r_node = node_diameter / 2 * 1.05
            x1o = x1 + ux * r_node
            y1o = y1 + uy * r_node
            x2o = x2 - ux * r_node
            y2o = y2 - uy * r_node
        else:
            x1o, y1o, x2o, y2o = x1, y1, x2, y2
        add_arrow(slide, x1o, y1o, x2o, y2o, c, width=w, head=head)

    # Draw nodes
    for idx, (x, y) in enumerate(positions):
        if perturbation_node is not None and idx == perturbation_node:
            # Perturbation source — slightly larger amber-filled node, white border.
            # No top-triangle marker (caused subtitle overlap); no inner dot
            # (caused donut artifact). Color + size do the visual work.
            big_d = node_diameter * 1.4
            add_oval(slide, x, y, big_d, ACCENT_AMBER, FG_PRIMARY, border_width=1.75)
        else:
            add_oval(slide, x, y, node_diameter, node_fill, node_border, border_width=1.0)


def build_slide2_causal(prs):
    slide = blank_slide(prs)
    add_title_block(slide,
                    "Causal Biological Intelligence",
                    "First learn structure. Then learn how signals flow.")

    # LEFT HALF — Topology
    left_cx = 3.4
    left_top = 1.7

    add_text(slide, 0.7, left_top, 5.5, 0.3,
             "TOPOLOGY LEARNING", size=12, bold=True, color=ACCENT_CYAN)
    add_text(slide, 0.7, left_top + 0.32, 5.5, 0.45,
             "Discover biological structure", size=20, bold=True, color=FG_PRIMARY)

    # Undirected network for topology — laid out lower with smaller radius
    # to clear subtitle text and leave room for bullet captions below.
    n_nodes = 8
    edges_topology = [
        (0, 1), (0, 2), (1, 3), (1, 4),
        (2, 4), (3, 5), (4, 5), (4, 6),
        (5, 7), (6, 7), (2, 6), (0, 7),
    ]
    network_cy = 3.85   # vertical center of network
    network_radius = 0.95
    _draw_network(
        slide,
        cx=left_cx,
        cy=network_cy,
        radius=network_radius,
        n_nodes=n_nodes,
        node_diameter=0.32,
        node_fill=ACCENT_CYAN,
        node_border=FG_PRIMARY,
        edges=edges_topology,
        edge_color=BORDER_SUBTLE,
        edge_width=1.25,
        directed=False,
        angle_offset=0.39,  # ~22.5° rotation so top node clears subtitle
    )

    bullet_top = network_cy + network_radius + 0.45
    add_multiline_text(
        slide, 0.9, bullet_top, 5.5, 0.75,
        [
            "Identify how biological components organize",
            "Build the latent map of cell biology",
        ],
        size=13, color=FG_PRIMARY, line_spacing=1.4,
    )

    # RIGHT HALF — Directional
    right_cx = 9.95
    right_top = 1.7

    add_text(slide, 7.25, right_top, 5.5, 0.3,
             "DIRECTIONAL CAUSAL LEARNING", size=12, bold=True, color=ACCENT_LAVENDER)
    add_text(slide, 7.25, right_top + 0.32, 5.5, 0.45,
             "Model perturbation flow", size=20, bold=True, color=FG_PRIMARY)

    # Directional network — same topology + arrow specs with varying weight
    arrow_specs = {
        (0, 1): (ACCENT_LAVENDER, 2.75, True),   # thick — primary outflow from perturbation
        (0, 2): (ACCENT_LAVENDER, 2.75, True),   # thick
        (1, 3): (ACCENT_LAVENDER, 1.75, True),   # medium
        (1, 4): (ACCENT_LAVENDER, 1.75, True),   # medium
        (2, 4): (ACCENT_LAVENDER, 1.75, True),   # medium
        (3, 5): (ACCENT_LAVENDER, 2.5, True),    # thick — secondary cascade
        (4, 5): (ACCENT_LAVENDER, 1.75, True),   # medium
        (4, 6): (BORDER_SUBTLE, 1.0, True),
        (5, 7): (BORDER_SUBTLE, 1.0, True),
        (6, 7): (BORDER_SUBTLE, 1.0, True),
        (2, 6): (BORDER_SUBTLE, 1.0, True),
        (0, 7): (BORDER_SUBTLE, 1.0, True),
    }
    _draw_network(
        slide,
        cx=right_cx,
        cy=network_cy,
        radius=network_radius,
        n_nodes=n_nodes,
        node_diameter=0.32,
        node_fill=ACCENT_LAVENDER,
        node_border=FG_PRIMARY,
        edges=edges_topology,
        edge_color=BORDER_SUBTLE,
        edge_width=1.25,
        directed=True,
        perturbation_node=0,
        arrow_specs=arrow_specs,
        angle_offset=0.39,  # match topology rotation; perturbation node sits upper-right
    )

    add_multiline_text(
        slide, 7.45, bullet_top, 5.5, 1.0,
        [
            "Infer directional signaling",
            "Estimate influence strength (edge bandwidth)",
            "Trace perturbation effects through network",
        ],
        size=13, color=FG_PRIMARY, line_spacing=1.4,
    )

    # BOTTOM STRIP — 3 validation cards (clear of bullet text above)
    cards_top = 6.3
    card_h = 0.85
    card_w = 3.85
    margin = 0.5
    gap = (13.333 - 2 * margin - 3 * card_w) / 2
    card_lefts = [margin + i * (card_w + gap) for i in range(3)]
    card_data = [
        ("Perturbation Validation", "Holds out perturbations, tests predictions"),
        ("Pathway Recovery", "Recovers known biological pathways from data"),
        ("Cross-State Consistency", "Consistent predictions across cell states"),
    ]
    for L, (title, desc) in zip(card_lefts, card_data):
        add_card(slide, L, cards_top, card_w, card_h,
                 fill_color=None, border_color=BORDER_SUBTLE, border_width=1.0, corner=0.08)
        # Green check icon — filled circle with a small triangle/check appearance
        icon_d = 0.22
        add_oval(slide, L + 0.28, cards_top + 0.27, icon_d,
                 ACCENT_GREEN, ACCENT_GREEN, border_width=0.5)
        add_text(slide, L + 0.55, cards_top + 0.13, card_w - 0.7, 0.32,
                 title, size=14, bold=True, color=FG_PRIMARY)
        add_text(slide, L + 0.55, cards_top + 0.43, card_w - 0.7, 0.42,
                 desc, size=11, color=FG_SECONDARY)

    # Footer note
    add_text(slide, 0.5, 7.18, 12.333, 0.3,
             '"AIVC learns both biological structure and how signals propagate through the system."',
             size=10, italic=True, color=FG_SECONDARY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 3 — Multimodal Encoder + Client Value
# ---------------------------------------------------------------------------
def build_slide3_encoder(prs):
    slide = blank_slide(prs)
    add_title_block(slide,
                    "Multimodal Encoder Architecture",
                    "Unifying complex multi-omics biology into actionable intelligence")

    # LEFT COLUMN — Architecture
    L_left = 0.55
    L_top = 1.75
    L_width = 6.1
    L_height = 5.0

    # Header row: ARCHITECTURE label (left) + legend (right) on same baseline
    add_text(slide, L_left, L_top, 2.5, 0.3,
             "ARCHITECTURE", size=12, bold=True, color=ACCENT_CYAN)
    add_text(slide, L_left + 2.5, L_top + 0.02, L_width - 2.5, 0.28,
             "✓ Today    ◆ Phase 1    ▲ Phase 2",
             size=10, italic=True, color=FG_SECONDARY, align=PP_ALIGN.RIGHT)

    # 5 input modality boxes
    in_top = L_top + 0.55
    in_h = 0.5
    margin_inner = 0.05
    n = 5
    gap_in = 0.12
    in_w = (L_width - 2 * margin_inner - 4 * gap_in) / n
    modality_labels = ["RNA", "Epigenetics", "Proteins", "Phosphoproteins", "VDJ"]
    modality_colors = [ACCENT_CYAN, ACCENT_CYAN, ACCENT_CYAN, ACCENT_LAVENDER, ACCENT_AMBER]
    in_lefts = [L_left + margin_inner + i * (in_w + gap_in) for i in range(n)]
    for x, label, color in zip(in_lefts, modality_labels, modality_colors):
        add_card(slide, x, in_top, in_w, in_h,
                 fill_color=None, border_color=color, border_width=1.5, corner=0.16)
        # Skip explicit icon; label centered for clarity
        add_text(slide, x, in_top + 0.06, in_w, in_h - 0.1,
                 label, size=10, bold=True, color=FG_PRIMARY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Central encoder block (legend moved to top — no gap eaten here)
    enc_top = in_top + in_h + 0.7
    enc_h = 0.9
    enc_w = L_width * 0.85
    enc_left = L_left + (L_width - enc_w) / 2
    enc_shape = add_card_translucent(
        slide, enc_left, enc_top, enc_w, enc_h,
        fill_color=ACCENT_LAVENDER, alpha_pct=22,
        border_color=ACCENT_LAVENDER, border_width=2.0, corner=0.10,
    )
    add_multiline_text(
        slide, enc_left, enc_top + 0.12, enc_w, enc_h - 0.15,
        [
            {"text": "Unified Multimodal Encoder", "size": 18, "bold": True, "color": FG_PRIMARY, "align": PP_ALIGN.CENTER},
            {"text": "256-dimensional latent representation", "size": 11, "italic": True, "color": FG_SECONDARY, "align": PP_ALIGN.CENTER},
        ],
        line_spacing=1.2,
    )

    # Arrows: from each input box down to encoder top
    for x, color in zip(in_lefts, modality_colors):
        x_center = x + in_w / 2
        add_arrow(slide, x_center, in_top + in_h + 0.05,
                  x_center, enc_top - 0.02,
                  color, width=1.0, head=True)

    # Output row
    out_top = enc_top + enc_h + 0.5
    out_h = 0.55
    n_out = 3
    gap_out = 0.18
    out_w = (L_width - 2 * margin_inner - 2 * gap_out) / n_out
    out_lefts = [L_left + margin_inner + i * (out_w + gap_out) for i in range(n_out)]
    out_labels = ["Biological State", "Perturbation Response", "Causal Inference"]
    for x, label in zip(out_lefts, out_labels):
        add_card(slide, x, out_top, out_w, out_h,
                 fill_color=None, border_color=ACCENT_CYAN, border_width=1.25, corner=0.16)
        add_text(slide, x, out_top + 0.08, out_w, out_h - 0.1,
                 label, size=11, bold=True, color=FG_PRIMARY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Arrows: from encoder bottom to each output
    enc_bottom_y = enc_top + enc_h
    for x in out_lefts:
        x_center = x + out_w / 2
        add_arrow(slide, x_center, enc_bottom_y + 0.02,
                  x_center, out_top - 0.02,
                  ACCENT_CYAN, width=1.0, head=True)

    # RIGHT COLUMN — Value cards 2x2
    R_left = 7.05
    R_top = 1.75
    R_width = 5.75
    R_height = 5.0

    add_text(slide, R_left, R_top, R_width, 0.3,
             "VALUE TO PARTNERS", size=12, bold=True, color=ACCENT_GREEN)

    cards = [
        ("Drug Response Prediction", "Predict combination efficacy"),
        ("Biomarker Discovery", "Identify patient stratification markers"),
        ("Target Prioritization", "Rank therapeutic targets by causal evidence"),
        ("Patient Stratification", "Match patients to optimal interventions"),
    ]
    # 2x2 grid
    grid_top = R_top + 0.55
    card_w = (R_width - 0.25) / 2
    card_h = (R_height - 0.6 - 0.25) / 2
    grid_gap = 0.25
    for idx, (title, desc) in enumerate(cards):
        row = idx // 2
        col = idx % 2
        x = R_left + col * (card_w + grid_gap)
        y = grid_top + row * (card_h + grid_gap)
        add_card(slide, x, y, card_w, card_h,
                 fill_color=None, border_color=BORDER_SUBTLE, border_width=1.2, corner=0.06)
        # Icon: small geometric — diamond shape (rotated square) at top-left
        icon_d = 0.28
        # Use diamond (MSO_SHAPE.DIAMOND)
        dia = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND,
            Inches(x + 0.25), Inches(y + 0.22),
            Inches(icon_d), Inches(icon_d),
        )
        dia.fill.solid()
        dia.fill.fore_color.rgb = ACCENT_GREEN
        dia.line.color.rgb = ACCENT_GREEN
        dia.line.width = Pt(0.5)
        dia.shadow.inherit = False
        # Title
        add_text(slide, x + 0.62, y + 0.22, card_w - 0.7, 0.4,
                 title, size=15, bold=True, color=FG_PRIMARY)
        # Description
        add_text(slide, x + 0.25, y + 0.7, card_w - 0.4, card_h - 0.8,
                 desc, size=12, color=FG_SECONDARY, line_spacing=1.3)

    # Bottom footer
    footer_top = 7.0
    add_hline(slide, 0.5, footer_top, 12.333, BORDER_SUBTLE, weight=0.75)
    add_text(slide, 0.5, footer_top + 0.12, 12.333, 0.3,
             "Designed to scale across datasets, perturbations, and therapeutic programs.",
             size=12, italic=True, color=FG_SECONDARY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 4 — Roadmap + Inflection Points
# ---------------------------------------------------------------------------
def build_slide4_roadmap(prs):
    slide = blank_slide(prs)
    add_title_block(slide,
                    "Roadmap & Key Inflection Points",
                    "Execution plan and platform value compounding")

    # TOP — Horizontal roadmap timeline
    axis_top = 1.95
    axis_left = 0.7
    axis_width = 11.933
    axis_y = axis_top + 1.05  # axis line in vertical middle of the roadmap area

    # Milestone data
    milestones = [
        ("Public Dataset Benchmarking", "2025 – Q2 2026",
         "Validated encoder · 73% cross-corpus", ACCENT_CYAN),
        ("Phase 1 Perturbation Learning", "Q3 2026",
         "Causal signal · BTK+JAK demo", ACCENT_LAVENDER),
        ("Phase 2: CRISPR + Multimodal", "2027",
         "5-modality scaling · 20–25 donors", ACCENT_AMBER),
        ("Scaled Data Generation", "2027 – 2028",
         "Cross-state reasoning", ACCENT_AMBER),
        ("Therapeutic Discovery Apps", "2028+",
         "Discovery enablement", ACCENT_GREEN),
    ]
    n_ms = len(milestones)

    # Axis line
    add_hline(slide, axis_left, axis_y, axis_width, BORDER_SUBTLE, weight=2.0)

    # Compute milestone x positions
    x_positions = [axis_left + (axis_width * (i + 0.5) / n_ms) for i in range(n_ms)]
    marker_d = 0.22

    for i, (title, timing, outcome, color) in enumerate(milestones):
        x = x_positions[i]
        # Filled circle marker on axis
        add_oval(slide, x, axis_y, marker_d, color, color, border_width=1.0)

        # Above axis: title + timing
        # Stagger alternate above/below to reduce vertical crowding? Spec says ABOVE label, BELOW description.
        # Title above
        label_w = (axis_width / n_ms) - 0.15
        add_text(slide, x - label_w / 2, axis_y - 0.95, label_w, 0.45,
                 title, size=12, bold=True, color=FG_PRIMARY,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM, line_spacing=1.1)
        # Timing above (smaller, color-coded)
        add_text(slide, x - label_w / 2, axis_y - 0.4, label_w, 0.3,
                 timing, size=10, bold=True, color=color,
                 align=PP_ALIGN.CENTER, line_spacing=1.1)

        # Below axis: outcome description
        add_text(slide, x - label_w / 2, axis_y + 0.3, label_w, 0.7,
                 outcome, size=10, italic=True, color=FG_SECONDARY,
                 align=PP_ALIGN.CENTER, line_spacing=1.2)

    # BOTTOM — Key Inflection Points (5 cards)
    add_text(slide, 0.5, 5.0, 12.333, 0.3,
             "KEY INFLECTION POINTS", size=12, bold=True, color=ACCENT_GREEN,
             align=PP_ALIGN.CENTER)

    cards_top = 5.45
    card_h = 1.3
    margin = 0.55
    n_cards = 5
    gap = 0.2
    card_w = (13.333 - 2 * margin - (n_cards - 1) * gap) / n_cards
    card_lefts = [margin + i * (card_w + gap) for i in range(n_cards)]
    inflection_cards = [
        ("Proprietary Multimodal Data", "Phase 1 lands", ACCENT_LAVENDER),
        ("Perturbation-Scale Expansion", "Phase 2 lands", ACCENT_AMBER),
        ("Causal Validation", "Architecture milestone", ACCENT_LAVENDER),
        ("Strategic Partnerships", "Pipeline programs begin", ACCENT_GREEN),
        ("Therapeutic Discovery", "Platform productization", ACCENT_GREEN),
    ]
    for L, (title, sub, color) in zip(card_lefts, inflection_cards):
        add_card(slide, L, cards_top, card_w, card_h,
                 fill_color=None, border_color=BORDER_SUBTLE, border_width=1.0, corner=0.08)
        # Top accent bar (colored stripe at top of card, small)
        stripe = add_rect(slide, L, cards_top, card_w, 0.06,
                          fill_color=color, border_color=None)
        # Icon: small diamond
        icon_d = 0.20
        dia = slide.shapes.add_shape(
            MSO_SHAPE.DIAMOND,
            Inches(L + card_w / 2 - icon_d / 2), Inches(cards_top + 0.18),
            Inches(icon_d), Inches(icon_d),
        )
        dia.fill.solid()
        dia.fill.fore_color.rgb = color
        dia.line.color.rgb = color
        dia.line.width = Pt(0.5)
        dia.shadow.inherit = False
        # Title
        add_text(slide, L + 0.1, cards_top + 0.46, card_w - 0.2, 0.55,
                 title, size=12, bold=True, color=FG_PRIMARY,
                 align=PP_ALIGN.CENTER, line_spacing=1.15)
        # Subline
        add_text(slide, L + 0.1, cards_top + 0.95, card_w - 0.2, 0.3,
                 sub, size=10, italic=True, color=FG_SECONDARY,
                 align=PP_ALIGN.CENTER, line_spacing=1.15)

    # Bottom footer note
    add_text(slide, 0.5, 7.05, 12.333, 0.3,
             '"Platform value compounds through proprietary data, causal learning, and scalable multimodal biological intelligence."',
             size=10, italic=True, color=FG_SECONDARY, align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide1_evolution(prs)
    build_slide2_causal(prs)
    build_slide3_encoder(prs)
    build_slide4_roadmap(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    size_kb = OUT_PATH.stat().st_size / 1024
    print(f"Saved: {OUT_PATH}  ({size_kb:.1f} KB)")


if __name__ == "__main__":
    main()
