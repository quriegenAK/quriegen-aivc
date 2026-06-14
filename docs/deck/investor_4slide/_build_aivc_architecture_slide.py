"""AIVC Architecture diagram — co-existing reasoning systems (Path A).

Architectural correction: GeneLink / Mechanica / PhysioMap are NOT sequential
phases. They are co-existing reasoning systems operating on ONE shared
biological state space, and they jointly compose into the Digital Cell /
Tissue Twin Layer.

Layout (top → bottom):

    ┌──────────────────────────┐
    │  Shared Latent           │   (cyan border, transparent fill)
    │  Biological State        │
    └─────────────┬────────────┘
                  │ (fan-out splitter)
   ┌──────────────┼──────────────┐
   ▼              ▼              ▼
┌────────┐   ┌────────┐    ┌────────┐
│GeneLink│   │Mechanica│   │PhysioMap│
│(data)  │   │(physics)│   │(pheno)  │
└────┬───┘   └────┬────┘   └────┬───┘
     │            │              │     (fan-in junction)
     └────────────┴──────┬───────┘
                         ▼
              ┌──────────────────────┐
              │  Digital Cell /      │  (green border, transparent fill)
              │  Tissue Twin Layer   │
              └──────────────────────┘

Visual notation: thin bordered transparent boxes for module boundaries.
This is a deliberate exception to v3's "no card containers" rule because
technical block diagrams REQUIRE module boundaries to communicate. The
"no cards" rule applies to content panels, not architectural notation.

Output: docs/deck/exports/aivc_architecture_v1.pptx (single slide)

Run:
  python3 docs/deck/investor_4slide/_build_aivc_architecture_slide.py
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Locked palette (matches v3 + Mechanica slide)
# ---------------------------------------------------------------------------
BG_DARK         = RGBColor(0x0A, 0x0E, 0x1A)
FG_PRIMARY      = RGBColor(0xFF, 0xFF, 0xFF)
FG_SECONDARY    = RGBColor(0xA0, 0xAF, 0xC8)
FG_MUTED        = RGBColor(0x60, 0x70, 0x88)
ACCENT_CYAN     = RGBColor(0x26, 0xDD, 0xF9)   # Shared State + GeneLink
ACCENT_LAVENDER = RGBColor(0x8B, 0x5C, 0xF6)   # Mechanica
ACCENT_AMBER    = RGBColor(0xF5, 0x9E, 0x0B)   # PhysioMap
ACCENT_GREEN    = RGBColor(0x4A, 0xDE, 0x80)   # Twin Layer
BORDER_SUBTLE   = RGBColor(0x2D, 0x3A, 0x57)

FONT = "Calibri"

REPO = Path(__file__).resolve().parents[3]
OUT_PATH = REPO / "docs" / "deck" / "exports" / "aivc_architecture_v1.pptx"


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------
def _zero_margins(tf):
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)


def _spc(run, hundredths_pt):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(hundredths_pt)))


def set_slide_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_text(slide, left, top, width, height, text,
             size=14, bold=False, italic=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, letter_spacing=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    _zero_margins(tf)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    if color is not None:
        run.font.color.rgb = color
    if letter_spacing is not None:
        _spc(run, letter_spacing)
    return box


def add_multiline(slide, left, top, width, height, lines,
                  size=13, color=None, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.4,
                  letter_spacing=None, bold=False, italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    _zero_margins(tf)
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        spec = item if isinstance(item, dict) else {"text": item}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", align)
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = spec["text"]
        run.font.size = Pt(spec.get("size", size))
        run.font.bold = spec.get("bold", bold)
        run.font.italic = spec.get("italic", italic)
        run.font.name = FONT
        col = spec.get("color", color)
        if col is not None:
            run.font.color.rgb = col
        ls = spec.get("letter_spacing", letter_spacing)
        if ls is not None:
            _spc(run, ls)
    return box


def add_block(slide, x, y, w, h, border_color, border_width=1.75, corner=0.06):
    """Thin bordered rounded rectangle, transparent fill — block-diagram convention."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.adjustments[0] = corner
    shape.fill.background()
    shape.line.color.rgb = border_color
    shape.line.width = Pt(border_width)
    shape.shadow.inherit = False
    return shape


def add_line(slide, x1, y1, x2, y2, color, width=1.5, head=False, head_size="med"):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if head:
        line_elem = conn.line._get_or_add_ln()
        for existing in line_elem.findall(qn("a:tailEnd")):
            line_elem.remove(existing)
        tail = etree.SubElement(line_elem, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", head_size)
        tail.set("len", head_size)
    return conn


def add_title_block(slide, title, subtitle):
    add_text(slide, 0.6, 0.55, 12.13, 0.7,
             title, size=44, bold=True, color=FG_PRIMARY,
             align=PP_ALIGN.LEFT, line_spacing=1.0)
    add_text(slide, 0.6, 1.25, 12.13, 0.4,
             subtitle, size=18, italic=True, color=FG_SECONDARY,
             align=PP_ALIGN.LEFT, line_spacing=1.1)
    add_line(slide, 0.6, 1.70, 12.733, 1.70, BORDER_SUBTLE, width=0.75)


def add_speaker_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# Layout constants
# ---------------------------------------------------------------------------
SLIDE_CENTER_X = 13.333 / 2  # 6.6665

# Shared State block (top)
SHARED_W, SHARED_H = 5.0, 0.90
SHARED_X = SLIDE_CENTER_X - SHARED_W / 2  # 4.1665
SHARED_Y = 1.95
SHARED_BOTTOM = SHARED_Y + SHARED_H       # 2.85

# Top splitter — fan-out from shared state to 3 modules
SPLITTER_Y = 3.15
DROP_TO_MODULE_END_Y = 3.40

# Module row
MODULE_Y = 3.50
MODULE_W = 3.65
MODULE_H = 1.75
MODULE_LEFTS = [0.60, 4.85, 9.05]
MODULE_CENTERS = [L + MODULE_W / 2 for L in MODULE_LEFTS]  # 2.425, 6.675, 10.875
MODULE_BOTTOM = MODULE_Y + MODULE_H       # 5.25

# Bottom fan-in junction
FANIN_Y = 5.55
FANIN_DROP_END_Y = 5.85

# Twin Layer (wide bottom block)
TWIN_W, TWIN_H = 9.50, 1.10
TWIN_X = SLIDE_CENTER_X - TWIN_W / 2      # 1.9165
TWIN_Y = 5.85
TWIN_BOTTOM = TWIN_Y + TWIN_H             # 6.95

# Footer
FOOTER_Y = 7.18


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------
def build_architecture_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    add_title_block(
        slide,
        title="AIVC Architecture",
        subtitle="Three co-existing reasoning systems over one shared biological state",
    )

    # ── Shared Latent Biological State block (top, cyan border) ───────────
    add_block(slide, SHARED_X, SHARED_Y, SHARED_W, SHARED_H,
              border_color=ACCENT_CYAN, border_width=2.0, corner=0.08)
    add_multiline(
        slide, SHARED_X, SHARED_Y + 0.13, SHARED_W, SHARED_H - 0.18,
        [
            {"text": "Shared Latent Biological State",
             "size": 20, "bold": True, "color": FG_PRIMARY,
             "align": PP_ALIGN.CENTER},
            {"text": "z ∈ ℝ²⁵⁶  ·  multi-omics encoder output",
             "size": 12, "italic": True, "color": ACCENT_CYAN,
             "align": PP_ALIGN.CENTER},
        ],
        line_spacing=1.25,
    )

    # ── Splitter: vertical from Shared down, horizontal fan-out, vertical drops ──
    # Vertical from shared bottom-center to splitter line
    add_line(slide, SLIDE_CENTER_X, SHARED_BOTTOM,
             SLIDE_CENTER_X, SPLITTER_Y,
             FG_SECONDARY, width=1.5)

    # Horizontal fan-out bar at SPLITTER_Y (spans the 3 module centers)
    add_line(slide, MODULE_CENTERS[0], SPLITTER_Y,
             MODULE_CENTERS[2], SPLITTER_Y,
             FG_SECONDARY, width=1.5)

    # 3 vertical drops to top of each module (with arrowheads)
    for cx in MODULE_CENTERS:
        add_line(slide, cx, SPLITTER_Y, cx, DROP_TO_MODULE_END_Y,
                 FG_SECONDARY, width=1.5, head=True, head_size="med")

    # ── Three module boxes ────────────────────────────────────────────────
    modules = [
        # (left, color, name, descriptor, body_lines)
        (MODULE_LEFTS[0], ACCENT_CYAN,
         "GeneLink", "data-driven",
         ["Learned biological relationships",
          "Causal structure from perturbation data"]),
        (MODULE_LEFTS[1], ACCENT_LAVENDER,
         "Mechanica", "physics-driven",
         ["Mechanistic simulation",
          "Dynamics · kinetics · resistance"]),
        (MODULE_LEFTS[2], ACCENT_AMBER,
         "PhysioMap", "phenotype-driven",
         ["Observable phenotype mapping",
          "Cell-state inference from morphology"]),
    ]
    for left, color, name, descriptor, body in modules:
        add_block(slide, left, MODULE_Y, MODULE_W, MODULE_H,
                  border_color=color, border_width=1.75, corner=0.06)
        cx = left + MODULE_W / 2
        # Title (centered)
        add_text(slide, left, MODULE_Y + 0.18, MODULE_W, 0.42,
                 name, size=24, bold=True, color=color,
                 align=PP_ALIGN.CENTER, line_spacing=1.0)
        # Descriptor (italic centered)
        add_text(slide, left, MODULE_Y + 0.65, MODULE_W, 0.32,
                 descriptor, size=13, italic=True, color=FG_SECONDARY,
                 align=PP_ALIGN.CENTER, line_spacing=1.0)
        # Body lines (centered, white)
        add_multiline(slide, left + 0.10, MODULE_Y + 1.05, MODULE_W - 0.20, 0.65,
                      [{"text": ln, "align": PP_ALIGN.CENTER} for ln in body],
                      size=12, color=FG_PRIMARY, line_spacing=1.35)

    # ── Bottom fan-in: 3 vertical drops from modules to junction, then single drop into twin ──
    for cx in MODULE_CENTERS:
        add_line(slide, cx, MODULE_BOTTOM, cx, FANIN_Y,
                 FG_SECONDARY, width=1.5)
    # Horizontal junction bar
    add_line(slide, MODULE_CENTERS[0], FANIN_Y,
             MODULE_CENTERS[2], FANIN_Y,
             FG_SECONDARY, width=1.5)
    # Single down-arrow from junction-center into twin layer top
    add_line(slide, SLIDE_CENTER_X, FANIN_Y,
             SLIDE_CENTER_X, FANIN_DROP_END_Y,
             FG_SECONDARY, width=1.75, head=True, head_size="med")

    # ── Digital Cell / Tissue Twin Layer block ────────────────────────────
    add_block(slide, TWIN_X, TWIN_Y, TWIN_W, TWIN_H,
              border_color=ACCENT_GREEN, border_width=2.0, corner=0.08)
    add_multiline(
        slide, TWIN_X, TWIN_Y + 0.18, TWIN_W, TWIN_H - 0.22,
        [
            {"text": "Digital Cell  /  Tissue Twin Layer",
             "size": 22, "bold": True, "color": FG_PRIMARY,
             "align": PP_ALIGN.CENTER},
            {"text": "Predictions  ·  Virtual perturbation experiments  ·  Digital twins",
             "size": 13, "italic": True, "color": ACCENT_GREEN,
             "align": PP_ALIGN.CENTER},
        ],
        line_spacing=1.30,
    )

    # ── Footer (key correction line) ──────────────────────────────────────
    add_text(slide, 0.6, FOOTER_Y, 12.13, 0.30,
             "These are co-existing reasoning systems, not sequential phases.",
             size=13, italic=True, color=FG_MUTED,
             align=PP_ALIGN.CENTER, line_spacing=1.1)

    add_speaker_notes(slide, _SPEAKER_NOTES)
    return slide


# ---------------------------------------------------------------------------
# Speaker notes
# ---------------------------------------------------------------------------
_SPEAKER_NOTES = """\
AIVC Architecture — Technical walk-through (90-120s)

CORRECTION FRAMING (10s):
"This is the architectural truth, not the temporal narrative. The
phase-based view we use with investors is a presentation simplification.
Internally, AIVC is three co-existing reasoning systems operating over
one shared biological state space."

TOP — SHARED STATE (15s):
"At the foundation: a single learned latent representation of biological
state — z, a 256-dimensional vector, produced by our multi-omics encoder.
Every reasoning system reads from and writes to the same state space.
This is the single most important architectural decision in the platform —
it's what makes the three modules ONE platform instead of three products."

THREE MODULES (40s):
"On top of that state, three complementary reasoning systems run in parallel.

GeneLink is the DATA-DRIVEN module — learns biological relationships and
causal structure from perturbation data. Answers 'what relationships exist.'

Mechanica is the PHYSICS-DRIVEN module — mechanistic simulation under
physics priors. Dynamics, kinetics, resistance evolution, combination
effects. Answers 'what happens when we perturb over time.'

PhysioMap is the PHENOTYPE-DRIVEN module — observable phenotype mapping,
cell-state inference from morphology and image-derived signals. Answers
'what is the observable cell state right now.'

Each operates on the same z but with a different reasoning bias. They
are complementary, not redundant. None subsumes the others."

BOTTOM — TWIN LAYER (20s):
"All three feed into the Digital Cell / Tissue Twin Layer. The twin layer
is the COMPOSITION of the three reasoning modes. A digital cell query —
'how does this patient's cell respond to this combination at 48 hours' —
routes through whichever module(s) have the best inference for that
specific question, with the others providing constraints or priors.

This is what powers predictions, virtual perturbation experiments, and
ultimately digital tissue twins."

EXPECTED DILIGENCE QUESTIONS:

Q: "Doesn't 3 modules over 1 state mean they're really one entangled model?"
A: They share the input space (z), not the parameters or the reasoning
architecture. Modules can be trained, validated, and shipped
independently. Composition happens at the query layer, not the training
layer.

Q: "Which module owns 'causality'?"
A: GeneLink learns causal structure from observed perturbations
(data-driven causality). Mechanica imposes mechanistic causality from
physics priors. They are two different routes to causal inference —
both legitimate, each with different generalization properties.

Q: "When does PhysioMap ship?"
A: PhysioMap is the imaging/morphology arm. Roadmap depends on imaging
data acquisition timeline. Mechanica is the immediate Phase 2 deliverable;
PhysioMap follows.

Q: "How do you avoid scope explosion across 3 modules?"
A: Single shared latent state space is the discipline. If any module
proposes a separate state representation, that's the trip-wire — we
collapse it before it ships.
"""


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_architecture_slide(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    size_kb = OUT_PATH.stat().st_size / 1024
    print(f"Saved: {OUT_PATH}  ({size_kb:.1f} KB)")


if __name__ == "__main__":
    main()
