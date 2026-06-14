"""AIVC Mechanica — board-level slide (Path A native shapes).

Spec: docs/deck/prompts/mechanica_slide_spec.md  (this file's accompanying spec)

Design language: v3 Apple-keynote minimal (typography + color + whitespace; no
card containers, no flattened images). Cross-slide visual language: same color
semantics + ●◆▲ symbol vocabulary + Calibri throughout + dark navy bg.

Output: docs/deck/exports/aivc_mechanica_v1.pptx (single slide)

Run:
  python3 docs/deck/investor_4slide/_build_mechanica_slide.py
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Locked palette (matches v3)
# ---------------------------------------------------------------------------
BG_DARK         = RGBColor(0x0A, 0x0E, 0x1A)
FG_PRIMARY      = RGBColor(0xFF, 0xFF, 0xFF)
FG_SECONDARY    = RGBColor(0xA0, 0xAF, 0xC8)
FG_MUTED        = RGBColor(0x60, 0x70, 0x88)
ACCENT_CYAN     = RGBColor(0x26, 0xDD, 0xF9)
ACCENT_LAVENDER = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_GREEN    = RGBColor(0x4A, 0xDE, 0x80)
BORDER_SUBTLE   = RGBColor(0x2D, 0x3A, 0x57)

FONT = "Calibri"

REPO = Path(__file__).resolve().parents[3]
OUT_PATH = REPO / "docs" / "deck" / "exports" / "aivc_mechanica_v1.pptx"


# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------
def _zero_margins(tf):
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)


def _spc(run, hundredths_pt):
    """Letter-spacing via XML `spc` attribute on rPr (Calibri compatible)."""
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(hundredths_pt)))


def set_slide_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_text(slide, left, top, width, height, text,
             size=14, bold=False, italic=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, letter_spacing=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    _zero_margins(tf)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    if color is not None:
        run.font.color.rgb = color
    if letter_spacing is not None:
        _spc(run, letter_spacing)
    return box


def add_multiline(slide, left, top, width, height, lines,
                  size=13, color=None, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.4,
                  letter_spacing=None, bold=False, italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                    Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    _zero_margins(tf)
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        spec = item if isinstance(item, dict) else {"text": item}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", align)
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = spec["text"]
        run.font.size = Pt(spec.get("size", size))
        run.font.bold = spec.get("bold", bold)
        run.font.italic = spec.get("italic", italic)
        run.font.name = FONT
        col = spec.get("color", color)
        if col is not None:
            run.font.color.rgb = col
        ls = spec.get("letter_spacing", letter_spacing)
        if ls is not None:
            _spc(run, ls)
    return box


def add_accent_rule(slide, x, y, length, color, thickness_in=0.05):
    """Solid filled rectangle used as a typographic accent above each phase column."""
    s = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y),
        Inches(length), Inches(thickness_in),
    )
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_divider(slide, x1, y, x2, color, width=0.75):
    """Subtle horizontal divider connector."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y), Inches(x2), Inches(y),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    return conn


def add_arrow(slide, x1, y1, x2, y2, color, width=1.5, head_size="med"):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1), Inches(y1), Inches(x2), Inches(y2),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    line_elem = conn.line._get_or_add_ln()
    for existing in line_elem.findall(qn("a:tailEnd")):
        line_elem.remove(existing)
    tail = etree.SubElement(line_elem, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", head_size)
    tail.set("len", head_size)
    return conn


def add_symbol(slide, cx, cy, symbol_type, color, size_in=0.34):
    """●  ◆  ▲  drawn as native shapes (circle / diamond / isosceles triangle).

    symbol_type: 'circle' | 'diamond' | 'triangle'
    """
    if symbol_type == "circle":
        mso = MSO_SHAPE.OVAL
    elif symbol_type == "diamond":
        mso = MSO_SHAPE.DIAMOND
    elif symbol_type == "triangle":
        mso = MSO_SHAPE.ISOSCELES_TRIANGLE
    else:
        raise ValueError(symbol_type)
    half = size_in / 2.0
    s = slide.shapes.add_shape(
        mso,
        Inches(cx - half), Inches(cy - half),
        Inches(size_in), Inches(size_in),
    )
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    s.shadow.inherit = False
    return s


def add_title_block(slide, title, subtitle):
    """Cross-slide title pattern: 44pt bold white + italic subtitle + thin divider."""
    add_text(slide, 0.6, 0.55, 12.13, 0.7,
             title, size=44, bold=True, color=FG_PRIMARY,
             align=PP_ALIGN.LEFT, line_spacing=1.0)
    add_text(slide, 0.6, 1.25, 12.13, 0.4,
             subtitle, size=18, italic=True, color=FG_SECONDARY,
             align=PP_ALIGN.LEFT, line_spacing=1.1)
    add_divider(slide, 0.6, 1.70, 12.733, BORDER_SUBTLE, width=0.75)


def add_speaker_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text


# ---------------------------------------------------------------------------
# Layout constants
# ---------------------------------------------------------------------------
# 3 phase columns (TODAY / EXPANSION / FUTURE)
PHASE_LEFTS  = [0.60, 4.85, 9.05]
PHASE_WIDTH  = 3.65
PHASE_CENTERS = [L + PHASE_WIDTH / 2 for L in PHASE_LEFTS]  # 2.425, 6.675, 10.875

# Vertical zones
HEADER_Y     = 1.95   # "TODAY" / "EXPANSION" / "FUTURE" small caps colored
ACCENT_Y     = 2.30   # accent rule (0.05" thick)
IDENTITY_Y   = 2.50   # 24pt bold colored (phase identity)
SYMBOL_Y     = 3.20   # phase symbol (centered in column)
QUESTION_Y   = 3.65   # italic question line, FG_SECONDARY
BODY_Y       = 4.05   # supporting body lines (15pt white)
BODY_H       = 1.20

# Divider between phase block and value block
VALUE_RULE_Y = 5.40

# Customer value block
VAL_HEADER_Y     = 5.55
VAL_SYMBOL_Y     = 6.05
VAL_NAME_Y       = 6.30
VAL_SUBTITLE_Y   = 6.78
PILLAR_CENTERS   = [1.85, 5.20, 8.55, 11.85]

# Vision closing line
VISION_Y = 7.18

# Phase connectors (horizontal arrows between phase columns at symbol baseline)
CONN_Y = SYMBOL_Y       # 3.20"
CONN_T_TO_E_X1, CONN_T_TO_E_X2 = 4.35, 4.80
CONN_E_TO_F_X1, CONN_E_TO_F_X2 = 8.55, 9.00


# ---------------------------------------------------------------------------
# Phase column renderer
# ---------------------------------------------------------------------------
def render_phase(slide, idx, color, header_label, identity, symbol_type,
                 question, body_lines):
    left = PHASE_LEFTS[idx]
    cx = PHASE_CENTERS[idx]

    # Header — uppercase small caps, letter-spaced, phase color
    add_text(slide, left, HEADER_Y, PHASE_WIDTH, 0.30,
             header_label, size=14, bold=True, color=color,
             align=PP_ALIGN.LEFT, line_spacing=1.0, letter_spacing=100)

    # Accent rule — 1.5" wide, left-aligned within column
    add_accent_rule(slide, left, ACCENT_Y, 1.50, color, thickness_in=0.05)

    # Identity statement — 24pt bold, phase color
    add_text(slide, left, IDENTITY_Y, PHASE_WIDTH, 0.55,
             identity, size=24, bold=True, color=color,
             align=PP_ALIGN.LEFT, line_spacing=1.0)

    # Symbol — centered in the column at SYMBOL_Y
    add_symbol(slide, cx, SYMBOL_Y, symbol_type, color, size_in=0.34)

    # Italic question line, FG_SECONDARY
    add_text(slide, left, QUESTION_Y, PHASE_WIDTH, 0.35,
             question, size=14, italic=True, color=FG_SECONDARY,
             align=PP_ALIGN.LEFT, line_spacing=1.15)

    # Body supporting lines — 15pt white, generous line spacing
    add_multiline(slide, left, BODY_Y, PHASE_WIDTH, BODY_H,
                  body_lines, size=15, color=FG_PRIMARY,
                  line_spacing=1.50)


# ---------------------------------------------------------------------------
# Slide build
# ---------------------------------------------------------------------------
def build_mechanica_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    add_title_block(
        slide,
        title="AIVC Mechanica",
        subtitle="The Mechanistic Intelligence Layer of AIVC",
    )

    # ── Phase columns ──────────────────────────────────────────────────────
    # TODAY: Multi-Omics + GeneLink (cyan, ● circle)
    render_phase(
        slide, 0, ACCENT_CYAN,
        header_label="TODAY",
        identity="Multi-Omics → GeneLink",
        symbol_type="circle",
        question="What biological relationships exist?",
        body_lines=[
            "Multi-omics encoder",
            "Learned causal relationships",
            "Validated on public + proprietary data",
        ],
    )

    # EXPANSION: Mechanica (lavender, ◆ diamond)
    render_phase(
        slide, 1, ACCENT_LAVENDER,
        header_label="EXPANSION",
        identity="AIVC Mechanica",
        symbol_type="diamond",
        question="What happens when the system is perturbed over time?",
        body_lines=[
            "Signaling dynamics",
            "Reaction kinetics",
            "Cell-state transitions",
            "Resistance evolution + combo effects",
        ],
    )

    # FUTURE: Digital Cells + Twins (amber, ▲ triangle)
    render_phase(
        slide, 2, ACCENT_AMBER,
        header_label="FUTURE",
        identity="Digital Cell Engine",
        symbol_type="triangle",
        question="What can be computed instead of measured?",
        body_lines=[
            "Digital Cells",
            "Virtual Perturbation Experiments",
            "Digital Tissue Twins",
        ],
    )

    # Phase connectors (horizontal arrows at symbol baseline)
    add_arrow(slide, CONN_T_TO_E_X1, CONN_Y, CONN_T_TO_E_X2, CONN_Y,
              FG_SECONDARY, width=1.5, head_size="med")
    add_arrow(slide, CONN_E_TO_F_X1, CONN_Y, CONN_E_TO_F_X2, CONN_Y,
              FG_SECONDARY, width=1.5, head_size="med")

    # ── Customer value block ───────────────────────────────────────────────
    add_divider(slide, 0.6, VALUE_RULE_Y, 12.733, BORDER_SUBTLE, width=0.75)
    add_text(slide, 0.6, VAL_HEADER_Y, 12.13, 0.32,
             "CUSTOMER VALUE", size=14, bold=True, color=ACCENT_GREEN,
             align=PP_ALIGN.LEFT, line_spacing=1.0, letter_spacing=100)

    pillars = [
        ("Better Target",         "Prioritization",  "Mechanism-aware ranking"),
        ("Earlier Resistance",    "Prediction",      "Simulate evolution in silico"),
        ("Rational Combination",  "Design",          "Kinetic-prior screening"),
        ("Reduced Experimental",  "Burden",          "Predict before you measure"),
    ]
    for cx, (line1, line2, subtitle) in zip(PILLAR_CENTERS, pillars):
        add_symbol(slide, cx, VAL_SYMBOL_Y, "diamond", ACCENT_GREEN, size_in=0.26)
        add_multiline(slide, cx - 1.6, VAL_NAME_Y, 3.2, 0.55,
                      [
                          {"text": line1, "align": PP_ALIGN.CENTER},
                          {"text": line2, "align": PP_ALIGN.CENTER},
                      ],
                      size=15, bold=True, color=FG_PRIMARY,
                      line_spacing=1.15)
        add_text(slide, cx - 1.7, VAL_SUBTITLE_Y, 3.4, 0.30,
                 subtitle, size=11, italic=True, color=FG_SECONDARY,
                 align=PP_ALIGN.CENTER, line_spacing=1.1)

    # ── Vision closing line ────────────────────────────────────────────────
    add_text(slide, 0.6, VISION_Y, 12.13, 0.30,
             "Mechanica becomes the computational engine powering digital cells "
             "and digital tissue twins.",
             size=12, italic=True, color=FG_MUTED,
             align=PP_ALIGN.CENTER, line_spacing=1.1)

    # ── Speaker notes (verbal walk-through for CEO / Ash) ──────────────────
    add_speaker_notes(slide, _SPEAKER_NOTES)
    return slide


# ---------------------------------------------------------------------------
# Speaker notes — 90-120 second verbal walk-through
# ---------------------------------------------------------------------------
_SPEAKER_NOTES = """\
AIVC Mechanica — Board-level walk-through (90-120 seconds)

OPENING (15s):
"AIVC has two complementary intelligence layers. GeneLink learns biological
relationships from data. Mechanica models how biology behaves under perturbation
over time. Together they are the foundation of our digital cell vision."

LEFT COLUMN — TODAY (20s):
"Today, our multi-omics encoder + GeneLink answer the question 'what biological
relationships exist?' We've validated 73% cross-corpus encoder generalization
and 57% four-way synergy classification on public perturbation data. This is
our learned foundation."

MIDDLE COLUMN — EXPANSION (30s):
"Mechanica is our next expansion. It adds the dynamic layer underneath GeneLink:
signaling dynamics, reaction kinetics, cell-state transitions, and crucially —
resistance evolution and combination therapy effects. Where GeneLink answers
'what relationships,' Mechanica answers 'what happens when we perturb over time.'
This is physics-informed simulation, constrained by mechanistic priors, operating
on the same latent state space as GeneLink. Shared encoder. Two reasoning modes."

RIGHT COLUMN — FUTURE (20s):
"Mechanica is the engine that makes digital cells real. Once we can simulate
perturbation dynamics with mechanistic confidence intervals, three capabilities
unlock: digital cells we can query in silico, virtual perturbation experiments
that replace bench work, and ultimately digital tissue twins for multi-cellular
context."

CUSTOMER VALUE (15s):
"For our pharma partners, this translates to four concrete value props: better
target prioritization through mechanism-aware ranking, earlier resistance
prediction via simulated evolution, rational combination design driven by
kinetic priors, and reduced experimental burden — predicting before measuring."

CLOSE (10s):
"Mechanica is not a competitor to GeneLink. It is the layer that lets GeneLink
extrapolate beyond what we've observed. Together they are AIVC's path from
'learned biology' to 'computable biology.'"

EXPECTED DILIGENCE QUESTIONS:
- "How is this different from CytoReason?" → CytoReason is literature-curated;
  Mechanica is learned-from-data with mechanistic priors. Our rate constants
  come from QuRIE-seq phospho timecourse — measured, not assumed.
- "When does Mechanica ship?" → Phase 2 deliverable; first integration with
  GeneLink Stage 3c causal architecture (Q1-Q2 2027).
- "What's the proof point on seed?" → Stage 3c resistance-prediction benchmark
  with Mechanica priors vs learned-only baseline. Quantitative delta.
- "Risk?" → Mechanistic ML for biology has a mixed empirical track record.
  We bench from day one against learned-only.
"""


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_mechanica_slide(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    size_kb = OUT_PATH.stat().st_size / 1024
    print(f"Saved: {OUT_PATH}  ({size_kb:.1f} KB)")


if __name__ == "__main__":
    main()
