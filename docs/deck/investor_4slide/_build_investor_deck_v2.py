"""Investor Slide 1 v2 — Pilot build, Path A native shapes.

Spec: docs/deck/prompts/slide1_v2_build_spec.md

Single-slide output validating the design system before extending to slides 2-4.
Design principles (non-negotiable):
  - 10-second comprehension test passes on Layer 1 (title + progression + headers)
  - Aggressive whitespace, no v1 ornamentation
  - Color discipline: cyan/lavender/amber per phase, dimmed for Phase 3
  - Native shapes only — no flattened images

Output: docs/deck/exports/aivc_investor_4slide_v2.pptx

Run:
  python3 docs/deck/investor_4slide/_build_investor_deck_v2.py
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Color palette (continued from v1, + DIMMED for Phase 3)
# ---------------------------------------------------------------------------
BG_DARK         = RGBColor(0x0A, 0x0E, 0x1A)
FG_PRIMARY      = RGBColor(0xFF, 0xFF, 0xFF)
FG_SECONDARY    = RGBColor(0xA0, 0xAF, 0xC8)
ACCENT_CYAN     = RGBColor(0x26, 0xDD, 0xF9)
ACCENT_LAVENDER = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
BORDER_SUBTLE   = RGBColor(0x2D, 0x3A, 0x57)
DIMMED          = RGBColor(0x60, 0x70, 0x88)

FONT = "Calibri"

REPO = Path(__file__).resolve().parents[3]
OUT_PATH = REPO / "docs" / "deck" / "exports" / "aivc_investor_4slide_v2.pptx"


# ---------------------------------------------------------------------------
# Low-level primitives
# ---------------------------------------------------------------------------
def set_slide_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _zero_textbox_margins(tf):
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)


def _set_run_letter_spacing(run, hundredths_pt):
    """Set character spacing on a run. PowerPoint stores this as `spc` in
    hundredths of a point on the rPr element. Positive = expanded.
    """
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(hundredths_pt)))


def _set_line_alpha(shape, alpha_percent):
    """Apply alpha to a shape's solid line. alpha_percent: 0-100."""
    ln = shape.line._get_or_add_ln()
    solidFill = ln.find(qn("a:solidFill"))
    if solidFill is None:
        return
    srgb = solidFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for existing in srgb.findall(qn("a:alpha")):
        srgb.remove(existing)
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", str(int(alpha_percent * 1000)))


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    size=14,
    bold=False,
    italic=False,
    color=None,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.15,
    letter_spacing=None,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    _zero_textbox_margins(tf)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    if color is not None:
        run.font.color.rgb = color
    if letter_spacing is not None:
        _set_run_letter_spacing(run, letter_spacing)
    return box


def add_multiline(
    slide,
    left,
    top,
    width,
    height,
    lines,
    size=13,
    color=None,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.4,
    letter_spacing=None,
    bold=False,
    italic=False,
):
    """Render a list of strings (or dicts) as paragraphs in one text box."""
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    _zero_textbox_margins(tf)
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        spec = item if isinstance(item, dict) else {"text": item}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", align)
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = spec["text"]
        run.font.size = Pt(spec.get("size", size))
        run.font.bold = spec.get("bold", bold)
        run.font.italic = spec.get("italic", italic)
        run.font.name = FONT
        col = spec.get("color", color)
        if col is not None:
            run.font.color.rgb = col
        ls = spec.get("letter_spacing", letter_spacing)
        if ls is not None:
            _set_run_letter_spacing(run, ls)
    return box


def add_card_outline(slide, left, top, width, height, border_color, border_alpha=100, border_width=1.5, corner=0.08):
    """Rounded rectangle, no fill (transparent), border only."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height),
    )
    shape.adjustments[0] = corner
    shape.fill.background()
    shape.line.color.rgb = border_color
    shape.line.width = Pt(border_width)
    if border_alpha < 100:
        _set_line_alpha(shape, border_alpha)
    shape.shadow.inherit = False
    return shape


def add_oval(slide, cx, cy, diameter, fill_color, border_color=None, border_width=1.0, dashed=False):
    r = diameter / 2.0
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - r), Inches(cy - r),
        Inches(diameter), Inches(diameter),
    )
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
        if dashed:
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shape.shadow.inherit = False
    return shape


def add_line(slide, x1, y, x2, color, width=1.5, dashed=False):
    """Horizontal line connector."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y), Inches(x2), Inches(y),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return conn


# ---------------------------------------------------------------------------
# Slide 1 v2
# ---------------------------------------------------------------------------
PANEL_TOP = 2.5
PANEL_HEIGHT = 4.0
PANEL_BOTTOM = PANEL_TOP + PANEL_HEIGHT
PANEL_PAD = 0.3
PANEL_LEFTS = [0.6, 4.74, 8.88]
PANEL_WIDTH = 3.84
PANEL_CENTERS = [L + PANEL_WIDTH / 2 for L in PANEL_LEFTS]  # 2.52, 6.66, 10.80

# Fixed inner-card y-coordinates (consistent across all 3 panels)
HEADER_Y = 2.85          # 18pt bold, multi-line (height auto)
IDENTITY_Y = 3.85        # 14pt regular, fixed across panels
SUPPORTING_Y = 4.20      # 13pt block of 3-4 lines
DIVIDER_Y = 5.55         # 1.5"-wide subtle line, panel color @30%
TAGLINE_Y = 5.80         # 12pt italic, centered, FG_SECONDARY

PROGRESSION_Y = 1.95     # horizontal line + dots y-center
DOT_DIAMETER = 0.16
DOT_X = [2.52, 6.66, 10.80, 12.20]


def _build_title_block(slide):
    add_text(
        slide, 0.6, 0.6, 12.13, 0.65,
        "AIVC Platform Evolution",
        size=40, bold=True, color=FG_PRIMARY,
        align=PP_ALIGN.LEFT, line_spacing=1.05,
    )
    add_text(
        slide, 0.6, 1.18, 12.13, 0.35,
        "From public benchmarking to scalable causal biological intelligence",
        size=16, italic=True, color=FG_SECONDARY,
        align=PP_ALIGN.LEFT, line_spacing=1.1,
    )


def _build_progression(slide):
    """Horizontal line + 4 dots + labels below each dot.

    Segments:
      dot1 → dot2 : solid BORDER_SUBTLE
      dot2 → dot3 : solid BORDER_SUBTLE
      dot3 → dot4 : DASHED BORDER_SUBTLE (discontinuity to future)
      dot4 → +0.4 : solid BORDER_SUBTLE (continuation indicator before arrow head)
    """
    # Solid segments connecting active phases
    add_line(slide, DOT_X[0], PROGRESSION_Y, DOT_X[1], BORDER_SUBTLE, width=1.5, dashed=False)
    add_line(slide, DOT_X[1], PROGRESSION_Y, DOT_X[2], BORDER_SUBTLE, width=1.5, dashed=False)
    # Dashed segment between Phase 2 and Phase 3 (signals temporal discontinuity)
    add_line(slide, DOT_X[2], PROGRESSION_Y, DOT_X[3], BORDER_SUBTLE, width=1.5, dashed=True)
    # Tiny continuation tail past Phase 3 dot to suggest "→ future"
    tail_end = DOT_X[3] + 0.35
    arrow_tail = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(DOT_X[3]), Inches(PROGRESSION_Y),
        Inches(tail_end), Inches(PROGRESSION_Y),
    )
    arrow_tail.line.color.rgb = BORDER_SUBTLE
    arrow_tail.line.width = Pt(1.0)
    arrow_tail.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    # Add arrowhead at tail
    line_elem = arrow_tail.line._get_or_add_ln()
    tail = etree.SubElement(line_elem, qn("a:tailEnd"))
    tail.set("type", "triangle")
    tail.set("w", "sm")
    tail.set("len", "sm")

    # Dots (drawn AFTER lines so they overlap on top cleanly)
    dot_specs = [
        (DOT_X[0], ACCENT_CYAN,     None,      False),  # NOW
        (DOT_X[1], ACCENT_LAVENDER, None,      False),  # PHASE 1
        (DOT_X[2], ACCENT_AMBER,    None,      False),  # PHASE 2
        (DOT_X[3], DIMMED,          FG_PRIMARY, True),  # Phase 3 (dashed white border)
    ]
    for x, fill, border, dashed in dot_specs:
        add_oval(slide, x, PROGRESSION_Y, DOT_DIAMETER,
                 fill_color=fill, border_color=border,
                 border_width=1.0, dashed=dashed)

    # Labels below dots
    label_top = PROGRESSION_Y + DOT_DIAMETER / 2 + 0.07  # 0.15" below dot edge
    label_height = 0.3
    # Each label is centered on the dot's x; we use a small txbox per label.
    def lbl(x_center, text, color, size=12, bold=True, italic=False, letter_spacing=100):
        w = 1.6
        add_text(slide, x_center - w / 2, label_top, w, label_height,
                 text, size=size, bold=bold, italic=italic, color=color,
                 align=PP_ALIGN.CENTER, line_spacing=1.0,
                 letter_spacing=letter_spacing)

    lbl(DOT_X[0], "NOW",      ACCENT_CYAN,     size=12, bold=True,  italic=False, letter_spacing=100)
    lbl(DOT_X[1], "PHASE 1",  ACCENT_LAVENDER, size=12, bold=True,  italic=False, letter_spacing=100)
    lbl(DOT_X[2], "PHASE 2",  ACCENT_AMBER,    size=12, bold=True,  italic=False, letter_spacing=100)
    lbl(DOT_X[3], "Phase 3",  FG_SECONDARY,    size=11, bold=False, italic=True,  letter_spacing=None)


def _build_panel(slide, panel_idx, color, header_lines, identity_line,
                 supporting_lines, tagline):
    left = PANEL_LEFTS[panel_idx]
    center_x = PANEL_CENTERS[panel_idx]

    # Card outline (transparent fill, panel-color border @ 60% alpha)
    add_card_outline(
        slide, left, PANEL_TOP, PANEL_WIDTH, PANEL_HEIGHT,
        border_color=color, border_alpha=60, border_width=1.5, corner=0.08,
    )

    # Inner padding
    inner_left = left + PANEL_PAD
    inner_width = PANEL_WIDTH - 2 * PANEL_PAD

    # Block 1 — Header (multi-line, panel color, 18pt bold, letter-spaced)
    header_height = 0.34 * len(header_lines) + 0.1  # rough fit
    add_multiline(
        slide, inner_left, HEADER_Y, inner_width, header_height,
        [{"text": ln, "letter_spacing": 50} for ln in header_lines],
        size=18, bold=True, color=color,
        align=PP_ALIGN.LEFT, line_spacing=1.05,
    )

    # Block 2 — Identity line (fixed y, single line, 14pt FG_PRIMARY)
    add_text(
        slide, inner_left, IDENTITY_Y, inner_width, 0.3,
        identity_line, size=14, color=FG_PRIMARY,
        align=PP_ALIGN.LEFT, line_spacing=1.2,
    )

    # Block 3 — Supporting details (fixed y, 13pt, generous line spacing)
    add_multiline(
        slide, inner_left, SUPPORTING_Y, inner_width, 1.5,
        supporting_lines,
        size=13, color=FG_PRIMARY, line_spacing=1.5,
    )

    # Block 4 — Divider (1.5" centered, panel color @ 30%)
    divider_w = 1.5
    divider_left = center_x - divider_w / 2
    div_conn = add_line(slide, divider_left, DIVIDER_Y, divider_left + divider_w,
                        color, width=1.0)
    _set_line_alpha(div_conn, 30)

    # Block 5 — Tagline (italic, centered)
    add_text(
        slide, left, TAGLINE_Y, PANEL_WIDTH, 0.3,
        tagline, size=12, italic=True, color=FG_SECONDARY,
        align=PP_ALIGN.CENTER, line_spacing=1.1,
    )


def _build_phase3_footer(slide):
    """Single line, italic. 'PHASE 3' prefix is bold. Arrow appended."""
    # Compose as multi-run in one text box for inline emphasis
    box = slide.shapes.add_textbox(Inches(0.6), Inches(6.8), Inches(12.13), Inches(0.35))
    tf = box.text_frame
    tf.word_wrap = True
    _zero_textbox_margins(tf)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.1

    def _run(text, bold=False, italic=False, color=FG_SECONDARY, size=13):
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = FONT
        r.font.color.rgb = color
        return r

    _run("PHASE 3", bold=True, italic=True, color=FG_SECONDARY, size=13)
    _run("  ───  ", italic=True, color=FG_SECONDARY, size=13)
    _run("Continuation at scale + therapeutic pipeline  →",
         italic=True, color=FG_SECONDARY, size=13)


def build_slide1_v2(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    _build_title_block(slide)
    _build_progression(slide)

    # Panel 1 — NOW
    _build_panel(
        slide, 0, ACCENT_CYAN,
        header_lines=["FOUNDATION", "& BENCHMARKING"],
        identity_line="Public multimodal datasets",
        supporting_lines=[
            "3 reference papers",
            "Pretrained encoder",
            "73% cross-corpus validation",
        ],
        tagline="Validated foundation",
    )

    # Panel 2 — PHASE 1
    _build_panel(
        slide, 1, ACCENT_LAVENDER,
        header_lines=["CONTROLLED", "PERTURBATION", "LEARNING"],
        identity_line="QuRIE-seq · proprietary multi-omics",
        supporting_lines=[
            "3 modalities (RNA · Protein · Phospho)",
            "5 donors · 5 timepoints",
            "5 stimuli · 10 inhibitors",
            "BTK + JAK headline demo",
        ],
        tagline="Causal learning, in motion",
    )

    # Panel 3 — PHASE 2
    _build_panel(
        slide, 2, ACCENT_AMBER,
        header_lines=["SCALABLE", "CAUSAL DISCOVERY"],
        identity_line="+ CRISPR + VDJ",
        supporting_lines=[
            "5 modalities (+ ATAC · VDJ)",
            "20–25 donors",
            "Soft + hard perturbations",
            "CRISPR screening library",
        ],
        tagline="Cross-state reasoning",
    )

    _build_phase3_footer(slide)
    return slide


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide1_v2(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    size_kb = OUT_PATH.stat().st_size / 1024
    print(f"Saved: {OUT_PATH}  ({size_kb:.1f} KB)")


if __name__ == "__main__":
    main()
