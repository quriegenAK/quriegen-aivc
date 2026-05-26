"""Investor 4-Slide Deck v3 — Apple-keynote minimal, Path A native shapes.

Spec: docs/deck/prompts/investor_4slide_v3_build_spec.md

Design direction (locked):
  - 3 levers do all the work: typography weight/scale, color saturation, whitespace
  - NO card containers anywhere across all 4 slides
  - Same color semantics + symbol vocabulary across slides
  - 5-tier hierarchy: title → visual narrative element → section labels → content → footnote

Output: docs/deck/exports/aivc_investor_4slide_v3.pptx

Run:
  python3 docs/deck/investor_4slide/_build_investor_deck_v3.py
"""
from __future__ import annotations

import math
from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Color palette (v3 = v2 + FG_MUTED for very subtle text)
# ---------------------------------------------------------------------------
BG_DARK         = RGBColor(0x0A, 0x0E, 0x1A)
FG_PRIMARY      = RGBColor(0xFF, 0xFF, 0xFF)
FG_SECONDARY    = RGBColor(0xA0, 0xAF, 0xC8)
FG_MUTED        = RGBColor(0x60, 0x70, 0x88)
ACCENT_CYAN     = RGBColor(0x26, 0xDD, 0xF9)
ACCENT_LAVENDER = RGBColor(0x8B, 0x5C, 0xF6)
ACCENT_AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
ACCENT_GREEN    = RGBColor(0x4A, 0xDE, 0x80)
BORDER_SUBTLE   = RGBColor(0x2D, 0x3A, 0x57)

FONT = "Calibri"

REPO = Path(__file__).resolve().parents[3]
OUT_PATH = REPO / "docs" / "deck" / "exports" / "aivc_investor_4slide_v3.pptx"


# ---------------------------------------------------------------------------
# Low-level primitives (carried + extended from v2)
# ---------------------------------------------------------------------------
def set_slide_bg(slide, rgb):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _zero_textbox_margins(tf):
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)


def _set_run_letter_spacing(run, hundredths_pt):
    rPr = run._r.get_or_add_rPr()
    rPr.set("spc", str(int(hundredths_pt)))


def _set_line_alpha(shape, alpha_percent):
    ln = shape.line._get_or_add_ln()
    solidFill = ln.find(qn("a:solidFill"))
    if solidFill is None:
        return
    srgb = solidFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for existing in srgb.findall(qn("a:alpha")):
        srgb.remove(existing)
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", str(int(alpha_percent * 1000)))


def _set_fill_alpha(shape, alpha_percent):
    spPr = shape.fill._xPr
    solidFill = spPr.find(qn("a:solidFill"))
    if solidFill is None:
        return
    srgb = solidFill.find(qn("a:srgbClr"))
    if srgb is None:
        return
    for existing in srgb.findall(qn("a:alpha")):
        srgb.remove(existing)
    alpha = etree.SubElement(srgb, qn("a:alpha"))
    alpha.set("val", str(int(alpha_percent * 1000)))


def add_text(slide, left, top, width, height, text,
             size=14, bold=False, italic=False, color=None,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             line_spacing=1.15, letter_spacing=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    _zero_textbox_margins(tf)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = FONT
    if color is not None:
        run.font.color.rgb = color
    if letter_spacing is not None:
        _set_run_letter_spacing(run, letter_spacing)
    return box


def add_multiline(slide, left, top, width, height, lines,
                  size=13, color=None, align=PP_ALIGN.LEFT,
                  anchor=MSO_ANCHOR.TOP, line_spacing=1.4,
                  letter_spacing=None, bold=False, italic=False):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    _zero_textbox_margins(tf)
    tf.vertical_anchor = anchor
    for i, item in enumerate(lines):
        spec = item if isinstance(item, dict) else {"text": item}
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = spec.get("align", align)
        if line_spacing:
            p.line_spacing = line_spacing
        run = p.add_run()
        run.text = spec["text"]
        run.font.size = Pt(spec.get("size", size))
        run.font.bold = spec.get("bold", bold)
        run.font.italic = spec.get("italic", italic)
        run.font.name = FONT
        col = spec.get("color", color)
        if col is not None:
            run.font.color.rgb = col
        ls = spec.get("letter_spacing", letter_spacing)
        if ls is not None:
            _set_run_letter_spacing(run, ls)
    return box


def add_inline_runs(slide, left, top, width, height, runs,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.1):
    """Single paragraph with mixed-style runs (for inline emphasis like the
    Phase 3 footer where 'Phase 3' is bold and the rest is italic).
    """
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    _zero_textbox_margins(tf)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for spec in runs:
        run = p.add_run()
        run.text = spec["text"]
        run.font.size = Pt(spec.get("size", 13))
        run.font.bold = spec.get("bold", False)
        run.font.italic = spec.get("italic", False)
        run.font.name = FONT
        col = spec.get("color")
        if col is not None:
            run.font.color.rgb = col
        ls = spec.get("letter_spacing")
        if ls is not None:
            _set_run_letter_spacing(run, ls)
    return box


def add_thick_accent_line(slide, x, y, length, color, thickness_in=0.06):
    """Solid filled rectangle used as a top accent line above each phase column
    (Slide 1) and as an emphasis rule elsewhere. Spec calls for ~0.18" thick
    but visually that's a heavy bar — use 0.06" for a crisp typographic accent.
    """
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y),
        Inches(length), Inches(thickness_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_modality_symbol(slide, cx, cy, symbol_type, color, size_in=0.32):
    """Geometric symbol used on Slide 3 modality row + Slide 4 inflection row.

    symbol_type: 'circle' | 'diamond' | 'triangle'
    Drawn at the given center coordinates.
    """
    if symbol_type == "circle":
        mso = MSO_SHAPE.OVAL
    elif symbol_type == "diamond":
        mso = MSO_SHAPE.DIAMOND
    elif symbol_type == "triangle":
        mso = MSO_SHAPE.ISOSCELES_TRIANGLE
    else:
        raise ValueError(symbol_type)
    half = size_in / 2.0
    shape = slide.shapes.add_shape(
        mso,
        Inches(cx - half), Inches(cy - half),
        Inches(size_in), Inches(size_in),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_oval(slide, cx, cy, diameter, fill_color, border_color=None,
             border_width=1.0, dashed=False):
    r = diameter / 2.0
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - r), Inches(cy - r),
        Inches(diameter), Inches(diameter),
    )
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    if border_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
        if dashed:
            shape.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    shape.shadow.inherit = False
    return shape


def add_line(slide, x1, y1, x2, y2, color, width=1.5, dashed=False,
             head=False, head_size="med", line_alpha=100):
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
        Inches(x2), Inches(y2),
    )
    conn.line.color.rgb = color
    conn.line.width = Pt(width)
    if dashed:
        conn.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    if head:
        line_elem = conn.line._get_or_add_ln()
        for existing in line_elem.findall(qn("a:tailEnd")):
            line_elem.remove(existing)
        tail = etree.SubElement(line_elem, qn("a:tailEnd"))
        tail.set("type", "triangle")
        tail.set("w", head_size)
        tail.set("len", head_size)
    if line_alpha < 100:
        _set_line_alpha(conn, line_alpha)
    return conn


def add_checkmark(slide, cx, cy, size_in, color):
    """Draw a check (✓) using two line segments at the given center.
    Native shapes only — no glyph dependency.
    """
    # Short stroke (down-right) + long stroke (up-right)
    s = size_in
    x_a1 = cx - s * 0.45
    y_a1 = cy
    x_a2 = cx - s * 0.10
    y_a2 = cy + s * 0.35
    x_b1 = x_a2
    y_b1 = y_a2
    x_b2 = cx + s * 0.50
    y_b2 = cy - s * 0.40
    for (x1, y1, x2, y2) in [(x_a1, y_a1, x_a2, y_a2), (x_b1, y_b1, x_b2, y_b2)]:
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(x1), Inches(y1), Inches(x2), Inches(y2),
        )
        conn.line.color.rgb = color
        conn.line.width = Pt(3.5)


def add_title_block(slide, title, subtitle):
    """Cross-slide visual language: 44pt bold white at top-left, italic subtitle below."""
    add_text(slide, 0.6, 0.55, 12.13, 0.7,
             title, size=44, bold=True, color=FG_PRIMARY,
             align=PP_ALIGN.LEFT, line_spacing=1.0)
    add_text(slide, 0.6, 1.25, 12.13, 0.4,
             subtitle, size=18, italic=True, color=FG_SECONDARY,
             align=PP_ALIGN.LEFT, line_spacing=1.1)


def blank_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s, BG_DARK)
    return s


# ===========================================================================
# SLIDE 1 — AIVC Platform Evolution
# ===========================================================================
S1_COL_LEFTS = [0.8, 4.9, 9.0]
S1_COL_WIDTH = 3.6
S1_COL_CENTERS = [L + S1_COL_WIDTH / 2 for L in S1_COL_LEFTS]  # 2.6, 6.7, 10.8
S1_PROG_Y = 2.10
S1_DOT_X = [S1_COL_CENTERS[0], S1_COL_CENTERS[1], S1_COL_CENTERS[2], 12.30]
S1_DOT_D = 0.22
S1_COL_TOP = 3.05
S1_ACCENT_W = 1.5
S1_ACCENT_THICK = 0.06
S1_HEADER_Y = S1_COL_TOP + 0.20
S1_IDENTITY_Y = S1_HEADER_Y + 0.55
S1_SUPPORTING_Y = S1_IDENTITY_Y + 1.05


def _s1_progression(slide):
    # Solid segments dot1↔dot2, dot2↔dot3
    add_line(slide, S1_DOT_X[0], S1_PROG_Y, S1_DOT_X[1], S1_PROG_Y,
             BORDER_SUBTLE, width=2.5)
    add_line(slide, S1_DOT_X[1], S1_PROG_Y, S1_DOT_X[2], S1_PROG_Y,
             BORDER_SUBTLE, width=2.5)
    # Dashed segment dot3 → dot4 (signals discontinuity to future)
    add_line(slide, S1_DOT_X[2], S1_PROG_Y, S1_DOT_X[3], S1_PROG_Y,
             BORDER_SUBTLE, width=2.5, dashed=True)
    # Continuation arrow past dot4
    add_line(slide, S1_DOT_X[3], S1_PROG_Y, S1_DOT_X[3] + 0.4, S1_PROG_Y,
             BORDER_SUBTLE, width=1.5, dashed=True, head=True, head_size="sm")

    # Active phase outer ring (PHASE 1 = dot index 1) — drawn BEFORE the dot
    # so it sits behind it visually.
    add_oval(slide, S1_DOT_X[1], S1_PROG_Y, 0.40,
             fill_color=None, border_color=ACCENT_LAVENDER,
             border_width=1.5)

    # Dots
    dot_specs = [
        (S1_DOT_X[0], ACCENT_CYAN,     None,       False),
        (S1_DOT_X[1], ACCENT_LAVENDER, None,       False),
        (S1_DOT_X[2], ACCENT_AMBER,    None,       False),
        (S1_DOT_X[3], FG_MUTED,        FG_PRIMARY, True),
    ]
    for x, fill, border, dashed in dot_specs:
        add_oval(slide, x, S1_PROG_Y, S1_DOT_D,
                 fill_color=fill, border_color=border,
                 border_width=1.0, dashed=dashed)

    # Labels below dots — PHASE 1 gets 16pt bold to emphasize active phase
    label_top = S1_PROG_Y + S1_DOT_D / 2 + 0.10
    label_w = 2.0
    label_h = 0.32

    def lbl(x_center, text, color, size, bold, italic, ls):
        add_text(slide, x_center - label_w / 2, label_top, label_w, label_h,
                 text, size=size, bold=bold, italic=italic, color=color,
                 align=PP_ALIGN.CENTER, line_spacing=1.0,
                 letter_spacing=ls)

    lbl(S1_DOT_X[0], "NOW",     ACCENT_CYAN,     14, True,  False, 100)
    lbl(S1_DOT_X[1], "PHASE 1", ACCENT_LAVENDER, 16, True,  False, 100)
    lbl(S1_DOT_X[2], "PHASE 2", ACCENT_AMBER,    14, True,  False, 100)
    lbl(S1_DOT_X[3], "Phase 3", FG_SECONDARY,    11, False, True,  None)


def _s1_column(slide, idx, color, header, identity_lines, supporting_lines):
    left = S1_COL_LEFTS[idx]
    # Top accent line — left-aligned within column
    add_thick_accent_line(slide, left, S1_COL_TOP, S1_ACCENT_W, color,
                          thickness_in=S1_ACCENT_THICK)

    # Phase header (uppercase, 22pt bold, panel color, letter-spaced)
    add_text(slide, left, S1_HEADER_Y, S1_COL_WIDTH, 0.40,
             header, size=22, bold=True, color=color,
             align=PP_ALIGN.LEFT, line_spacing=1.0,
             letter_spacing=100)

    # Identity statement (multi-line short — 18pt regular white)
    add_multiline(slide, left, S1_IDENTITY_Y, S1_COL_WIDTH, 1.0,
                  identity_lines,
                  size=18, color=FG_PRIMARY, align=PP_ALIGN.LEFT,
                  line_spacing=1.20)

    # Supporting lines (15pt regular white, generous line spacing)
    add_multiline(slide, left, S1_SUPPORTING_Y, S1_COL_WIDTH, 2.8,
                  supporting_lines,
                  size=15, color=FG_PRIMARY, align=PP_ALIGN.LEFT,
                  line_spacing=1.45)


def build_slide1(prs):
    slide = blank_slide(prs)
    add_title_block(
        slide,
        "AIVC Platform Evolution",
        "From public benchmarking to scalable causal biological intelligence",
    )
    _s1_progression(slide)

    # Column 1 — NOW (cyan)
    _s1_column(
        slide, 0, ACCENT_CYAN,
        header="NOW",
        identity_lines=["Public", "benchmarking"],
        supporting_lines=[
            "3 reference papers",
            "Pretrained encoder",
            "73% cross-corpus validation",
        ],
    )

    # Column 2 — PHASE 1 (lavender, active focus)
    _s1_column(
        slide, 1, ACCENT_LAVENDER,
        header="PHASE 1",
        identity_lines=["Proprietary", "perturbation", "learning"],
        supporting_lines=[
            "QuRIE-seq multi-omics",
            "3 modalities",
            "5 donors · 5 timepoints",
            "5 stimuli · 10 inhibitors",
            {"text": "BTK + JAK demo", "color": ACCENT_LAVENDER, "bold": True},
        ],
    )

    # Column 3 — PHASE 2 (amber)
    _s1_column(
        slide, 2, ACCENT_AMBER,
        header="PHASE 2",
        identity_lines=["Scaled", "causal", "discovery"],
        supporting_lines=[
            "+ CRISPR + VDJ",
            "5 modalities",
            "20–25 donors",
            "Soft + hard perturbations",
            {"text": "CRISPR screening library", "color": ACCENT_AMBER, "bold": True},
        ],
    )

    # Phase 3 footer — single inline line, centered
    add_inline_runs(
        slide, 0.6, 6.85, 12.13, 0.32,
        runs=[
            {"text": "Phase 3", "bold": True,  "italic": True, "color": FG_MUTED, "size": 13},
            {"text": " — Continuation at scale + therapeutic pipeline  →",
             "italic": True, "color": FG_MUTED, "size": 13},
        ],
        align=PP_ALIGN.CENTER,
    )


# ===========================================================================
# SLIDE 2 — Causal Biological Intelligence
# ===========================================================================
# 9-node hex-cluster layout (relative offsets in inches from network center).
# Slightly organic — not perfectly hexagonal.
S2_NODE_OFFSETS = [
    ( 0.00, -1.10),   # 0 — top  (source on right network)
    ( 1.00, -0.55),   # 1 — top-right
    ( 1.10,  0.50),   # 2 — bottom-right
    ( 0.40,  1.10),   # 3 — bottom-right-low
    (-0.40,  1.15),   # 4 — bottom-left-low
    (-1.05,  0.55),   # 5 — bottom-left
    (-1.10, -0.55),   # 6 — top-left
    (-0.05,  0.10),   # 7 — center
    ( 0.55, -0.30),   # 8 — mid-right
]
S2_EDGES = [
    (0, 1), (0, 6), (0, 7),
    (1, 2), (1, 7), (1, 8),
    (2, 3), (2, 7),
    (3, 4), (3, 7),
    (4, 5), (4, 7),
    (5, 6), (5, 7),
    (6, 7), (7, 8),
]


def _s2_network(slide, cx, cy, node_color, directional=False, source_idx=0):
    # Compute absolute node positions
    nodes = [(cx + dx, cy + dy) for (dx, dy) in S2_NODE_OFFSETS]
    node_d = 0.32

    # For directional, classify edges into thick / medium / thin
    arrow_specs = {}
    if directional:
        thick_edges = [(0, 1), (0, 7), (0, 6)]
        medium_edges = [(1, 2), (5, 6), (7, 8), (3, 7)]
        for e in S2_EDGES:
            if e in thick_edges:
                arrow_specs[e] = (ACCENT_LAVENDER, 3.0, 100)
            elif e in medium_edges:
                arrow_specs[e] = (ACCENT_LAVENDER, 2.0, 70)
            else:
                arrow_specs[e] = (BORDER_SUBTLE, 1.0, 100)

    # Draw edges first (so nodes overlap on top)
    for (i, j) in S2_EDGES:
        x1, y1 = nodes[i]
        x2, y2 = nodes[j]
        if directional:
            c, w, alpha = arrow_specs[(i, j)]
            # Offset endpoints to clear node circles
            dx, dy = x2 - x1, y2 - y1
            dist = (dx * dx + dy * dy) ** 0.5
            ux, uy = (dx / dist, dy / dist) if dist > 0 else (0, 0)
            r_node = node_d / 2 * 1.05
            add_line(slide,
                     x1 + ux * r_node, y1 + uy * r_node,
                     x2 - ux * r_node, y2 - uy * r_node,
                     color=c, width=w, head=True, head_size="sm",
                     line_alpha=alpha)
        else:
            add_line(slide, x1, y1, x2, y2,
                     color=BORDER_SUBTLE, width=1.5)

    # Source node ring (drawn BEFORE the source node so it sits behind)
    if directional and source_idx is not None:
        sx, sy = nodes[source_idx]
        # Outer concentric ring — 0.46" diameter (vs 0.32" node), lavender 1.25pt
        add_oval(slide, sx, sy, 0.50,
                 fill_color=None, border_color=ACCENT_LAVENDER,
                 border_width=1.25)

    # Draw nodes
    for idx, (nx, ny) in enumerate(nodes):
        if directional and idx == source_idx:
            # Source = brighter / slightly larger
            add_oval(slide, nx, ny, node_d * 1.08,
                     fill_color=ACCENT_LAVENDER, border_color=FG_PRIMARY,
                     border_width=1.0)
        else:
            add_oval(slide, nx, ny, node_d,
                     fill_color=node_color, border_color=None)


def build_slide2(prs):
    slide = blank_slide(prs)
    add_title_block(
        slide,
        "Causal Biological Intelligence",
        "First learn structure. Then learn how signals flow.",
    )

    left_cx, right_cx = 2.8, 10.4
    network_cy = 3.35

    # Left network — undirected topology
    _s2_network(slide, left_cx, network_cy,
                node_color=ACCENT_CYAN, directional=False)

    # Right network — directional, source = node 0 (top)
    _s2_network(slide, right_cx, network_cy,
                node_color=ACCENT_LAVENDER, directional=True, source_idx=0)

    # Connecting arrow between networks (narrative "first this, then this")
    # Use a thick rectangle + triangular head shape for a deliberate arrow.
    arrow_left, arrow_right, arrow_y = 4.55, 8.55, network_cy
    add_line(slide, arrow_left, arrow_y, arrow_right - 0.25, arrow_y,
             color=FG_SECONDARY, width=4.0)
    # Arrowhead triangle
    head = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_TRIANGLE,
        Inches(arrow_right - 0.25), Inches(arrow_y - 0.13),
        Inches(0.28), Inches(0.26),
    )
    # The default RIGHT_TRIANGLE has right-angle bottom-left; rotate so the
    # hypotenuse points right (we want a forward-pointing arrowhead). The
    # simpler reliable approach: use ISOCELES_TRIANGLE rotated 90° clockwise.
    head_alt = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(arrow_right - 0.30), Inches(arrow_y - 0.18),
        Inches(0.32), Inches(0.36),
    )
    head_alt.rotation = 90
    head_alt.fill.solid()
    head_alt.fill.fore_color.rgb = FG_SECONDARY
    head_alt.line.fill.background()
    head_alt.shadow.inherit = False
    # Remove the placeholder RIGHT_TRIANGLE we added first
    sp = head._element
    sp.getparent().remove(sp)

    # Network labels (5.05"-5.55")
    add_text(slide, left_cx - 2.0, 4.95, 4.0, 0.4,
             "TOPOLOGY", size=22, bold=True, color=ACCENT_CYAN,
             align=PP_ALIGN.CENTER, letter_spacing=100, line_spacing=1.0)
    add_text(slide, left_cx - 2.5, 5.40, 5.0, 0.4,
             "Discover structure", size=18, color=FG_PRIMARY,
             align=PP_ALIGN.CENTER, line_spacing=1.1)

    add_text(slide, right_cx - 2.0, 4.95, 4.0, 0.4,
             "DIRECTIONAL", size=22, bold=True, color=ACCENT_LAVENDER,
             align=PP_ALIGN.CENTER, letter_spacing=100, line_spacing=1.0)
    add_text(slide, right_cx - 2.5, 5.40, 5.0, 0.4,
             "Model perturbation flow", size=18, color=FG_PRIMARY,
             align=PP_ALIGN.CENTER, line_spacing=1.1)

    # Validation row at bottom (5.9"-6.9") — 3 text columns with checkmark
    val_cols = [
        (2.3, "Perturbation validation", "Held-out perturbations match predictions"),
        (6.665, "Pathway recovery",        "Recovers known biological pathways"),
        (11.0, "Cross-state consistency", "Stable across cell states"),
    ]
    for cx, title, desc in val_cols:
        add_checkmark(slide, cx, 6.05, 0.32, ACCENT_GREEN)
        add_text(slide, cx - 2.0, 6.30, 4.0, 0.35,
                 title, size=15, bold=True, color=FG_PRIMARY,
                 align=PP_ALIGN.CENTER, line_spacing=1.1)
        add_text(slide, cx - 2.25, 6.62, 4.5, 0.32,
                 desc, size=12, color=FG_SECONDARY,
                 align=PP_ALIGN.CENTER, line_spacing=1.15)


# ===========================================================================
# SLIDE 3 — Multimodal Encoder + Value
# ===========================================================================
S3_MODALITIES = [
    # (label, x_center, color, symbol)
    ("RNA",      2.0,  ACCENT_CYAN,     "circle"),
    ("ATAC",     4.4,  ACCENT_CYAN,     "circle"),
    ("Protein",  6.665, ACCENT_CYAN,    "circle"),
    ("Phospho",  9.0,  ACCENT_LAVENDER, "diamond"),
    ("VDJ",      11.3, ACCENT_AMBER,    "triangle"),
]
S3_VALUES = [
    # (x_center, name (one or two strings), subtitle)
    (1.95,  ["Drug response", "prediction"],   "Predict combinations"),
    (5.30,  ["Biomarker",     "discovery"],    "Stratify patients"),
    (8.30,  ["Target",        "prioritization"], "Rank by causal evidence"),
    (11.55, ["Patient",       "stratification"], "Match to interventions"),
]


def build_slide3(prs):
    slide = blank_slide(prs)
    add_title_block(
        slide,
        "Multimodal Encoder",
        "Unifying multi-omics biology into actionable intelligence",
    )

    # MODALITY ROW (2.0-2.9)
    name_y = 2.05
    symbol_y = 2.65
    for label, x, color, symbol in S3_MODALITIES:
        add_text(slide, x - 1.5, name_y, 3.0, 0.35,
                 label, size=18, bold=True, color=color,
                 align=PP_ALIGN.CENTER, line_spacing=1.1)
        add_modality_symbol(slide, x, symbol_y, symbol, color, size_in=0.34)

    # Convergence arrows from each modality DOWN to encoder bar top
    encoder_top = 3.75
    encoder_left, encoder_w, encoder_h = 1.5, 10.333, 0.95
    for label, x, color, symbol in S3_MODALITIES:
        add_line(slide, x, symbol_y + 0.22, x, encoder_top - 0.05,
                 color=FG_SECONDARY, width=1.0, head=True, head_size="sm")

    # ENCODER BAR (3.75-4.70)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(encoder_left), Inches(encoder_top),
        Inches(encoder_w), Inches(encoder_h),
    )
    bar.adjustments[0] = 0.06
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_LAVENDER
    _set_fill_alpha(bar, 15)
    bar.line.color.rgb = ACCENT_LAVENDER
    bar.line.width = Pt(1.75)
    bar.shadow.inherit = False
    add_multiline(slide, encoder_left, encoder_top + 0.12, encoder_w,
                  encoder_h - 0.15,
                  [
                      {"text": "UNIFIED ENCODER", "size": 24, "bold": True,
                       "color": FG_PRIMARY, "align": PP_ALIGN.CENTER,
                       "letter_spacing": 100},
                      {"text": "256-D latent representation",
                       "size": 14, "italic": True, "color": FG_SECONDARY,
                       "align": PP_ALIGN.CENTER},
                  ],
                  line_spacing=1.2)

    # Single thick down-arrow from encoder center to value row
    encoder_bottom = encoder_top + encoder_h
    add_line(slide, 6.665, encoder_bottom + 0.05,
             6.665, encoder_bottom + 0.55,
             color=ACCENT_LAVENDER, width=4.0, head=True, head_size="lg")

    # VALUE ROW (5.4-6.3)
    val_name_y = 5.40
    val_sub_y = 6.05
    for cx, name_lines, subtitle in S3_VALUES:
        add_multiline(slide, cx - 1.5, val_name_y, 3.0, 0.65,
                      [{"text": ln, "align": PP_ALIGN.CENTER} for ln in name_lines],
                      size=16, bold=True, color=FG_PRIMARY,
                      line_spacing=1.1)
        add_text(slide, cx - 1.6, val_sub_y, 3.2, 0.32,
                 subtitle, size=12, italic=True, color=FG_SECONDARY,
                 align=PP_ALIGN.CENTER, line_spacing=1.1)

    # LEGEND STRIP (6.7-7.0) — inline mixed-color runs
    add_inline_runs(
        slide, 0.6, 6.85, 12.13, 0.32,
        runs=[
            {"text": "●", "color": ACCENT_CYAN,     "size": 14, "bold": True},
            {"text": "  Today      ", "color": FG_SECONDARY, "size": 11},
            {"text": "◆", "color": ACCENT_LAVENDER, "size": 14, "bold": True},
            {"text": "  Phase 1      ", "color": FG_SECONDARY, "size": 11},
            {"text": "▲", "color": ACCENT_AMBER,    "size": 14, "bold": True},
            {"text": "  Phase 2", "color": FG_SECONDARY, "size": 11},
        ],
        align=PP_ALIGN.CENTER,
    )


# ===========================================================================
# SLIDE 4 — Roadmap & Inflection Points
# ===========================================================================
S4_TIMELINE_Y = 2.55
S4_DOT_D = 0.22
S4_MILESTONES = [
    # (x, year, outcome_lines, color)
    (1.5,  "2025",    ["Public benchmark"],               ACCENT_CYAN),
    (4.0,  "Q3'26",   ["Phase 1 data", "lands"],          ACCENT_LAVENDER),
    (6.665,"2027",    ["Phase 2 CRISPR", "+ VDJ"],        ACCENT_AMBER),
    (9.2,  "2027-28", ["Scaled data", "generation"],      ACCENT_AMBER),
    (11.8, "2028+",   ["Therapeutic", "discovery"],       ACCENT_GREEN),
]
S4_INFLECTIONS = [
    # (x, symbol, color, name_lines)
    (2.0,  "circle",   ACCENT_CYAN,     ["Proprietary", "multimodal data"]),
    (5.0,  "diamond",  ACCENT_LAVENDER, ["Perturbation", "scale expansion"]),
    (8.0,  "diamond",  ACCENT_LAVENDER, ["Causal", "validation"]),
    (11.0, "triangle", ACCENT_AMBER,    ["Therapeutic discovery", "enablement"]),
]


def build_slide4(prs):
    slide = blank_slide(prs)
    add_title_block(
        slide,
        "Roadmap",
        "Execution plan and value compounding",
    )

    # TIMELINE LINE (y=2.55, 1.5 → 11.8)
    add_line(slide, 1.5, S4_TIMELINE_Y, 11.8, S4_TIMELINE_Y,
             color=BORDER_SUBTLE, width=2.5)

    # DOTS + YEAR + OUTCOME
    for x, year, outcome, color in S4_MILESTONES:
        # Dot
        add_oval(slide, x, S4_TIMELINE_Y, S4_DOT_D,
                 fill_color=color, border_color=None)
        # Year label below (centered on dot)
        add_text(slide, x - 1.2, S4_TIMELINE_Y + 0.20, 2.4, 0.32,
                 year, size=14, bold=True, color=color,
                 align=PP_ALIGN.CENTER, line_spacing=1.0,
                 letter_spacing=50)
        # Outcome label (multi-line)
        add_multiline(slide, x - 1.4, S4_TIMELINE_Y + 0.65, 2.8, 0.85,
                      [{"text": ln, "align": PP_ALIGN.CENTER} for ln in outcome],
                      size=15, color=FG_PRIMARY, line_spacing=1.2)

    # INFLECTION POINTS SECTION (header @ y=5.0)
    add_text(slide, 0.8, 5.00, 12.13, 0.40,
             "Inflection points", size=22, bold=True, color=ACCENT_GREEN,
             align=PP_ALIGN.LEFT, line_spacing=1.0,
             letter_spacing=50)

    # 4 columns @ y=5.6 symbol top
    symbol_y = 5.85
    name_y = 6.20
    for x, symbol, color, name_lines in S4_INFLECTIONS:
        add_modality_symbol(slide, x, symbol_y, symbol, color, size_in=0.32)
        add_multiline(slide, x - 1.7, name_y, 3.4, 0.85,
                      [{"text": ln, "align": PP_ALIGN.CENTER} for ln in name_lines],
                      size=16, bold=True, color=FG_PRIMARY, line_spacing=1.2)


# ===========================================================================
# Main
# ===========================================================================
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide1(prs)
    build_slide2(prs)
    build_slide3(prs)
    build_slide4(prs)

    OUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT_PATH)
    size_kb = OUT_PATH.stat().st_size / 1024
    print(f"Saved: {OUT_PATH}  ({size_kb:.1f} KB)")


if __name__ == "__main__":
    main()
